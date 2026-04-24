"""
PptxScreenshotExporter — Render each slide of a .pptx file to a PNG.

Two backends are supported, selectable via the ``backend`` constructor argument:

``powerpoint``
    Uses the PowerPoint COM automation API (Windows only).  Requires pywin32.
    Provides the highest-quality rendering (identical to what PowerPoint shows).
    ``slide.Export(path, "PNG", width, height)`` per slide.

``libreoffice``
    Uses LibreOffice in headless mode via subprocess.  Cross-platform (Linux /
    Ubuntu).  Requires LibreOffice >= 7.0 on PATH.  Converts the whole .pptx
    in one invocation, then renames the numbered output files to the
    ``slide-001-Title.png`` convention.

``auto`` (default)
    Tries the ``powerpoint`` backend first (Windows + pywin32 available).
    Falls back to ``libreoffice`` if PowerPoint is not available.
"""

from __future__ import annotations

import logging
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

logger = logging.getLogger(__name__)

# Points-to-pixels at standard 96 DPI (PowerPoint stores dimensions in EMU / points)
_PT_TO_PX = 96.0 / 72.0


def _safe_name(title: str) -> str:
    """Sanitise a slide title for use in a filename."""
    return re.sub(r"[^\w]+", "_", title).strip("_") or "Slide"


def _slide_titles_from_pptx(pptx_path: Path) -> list[str]:
    """Extract slide titles using python-pptx (already a dependency)."""
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
        prs = Presentation(str(pptx_path))
        titles: list[str] = []
        for i, slide in enumerate(prs.slides):
            title = None
            # Try the slide name first (set by our exporter), then placeholder title
            if slide.name and not slide.name.startswith("Slide "):
                title = slide.name
            else:
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.shape_type == 13:  # title placeholder type
                        title = shape.text_frame.text.strip()
                        break
                    if hasattr(shape, "placeholder_format") and shape.placeholder_format is not None:
                        if shape.placeholder_format.idx == 0:  # title placeholder
                            title = shape.text_frame.text.strip() if shape.has_text_frame else None
                            break
            titles.append(title or f"Slide {i + 1}")
        return titles
    except Exception as exc:  # noqa: BLE001
        logger.warning("Could not read slide titles from PPTX: %s", exc)
        return []


def _canvas_size_from_pptx(pptx_path: Path) -> tuple[int, int]:
    """Read slide canvas dimensions (px) from the PPTX file itself."""
    try:
        from pptx import Presentation
        from pptx.util import Emu
        prs = Presentation(str(pptx_path))
        # slide_width / slide_height are in EMU (914400 EMU per inch, 96 px per inch)
        w_px = int(prs.slide_width / 914400 * 96)
        h_px = int(prs.slide_height / 914400 * 96)
        return w_px, h_px
    except Exception as exc:  # noqa: BLE001
        logger.warning("Could not read canvas size from PPTX (%s); using 1280×720", exc)
        return 1280, 720


# ── PowerPoint backend ───────────────────────────────────────────────────────

def _export_powerpoint(
    pptx_path: Path,
    out_dir: Path,
    scale: int,
    slide_titles: list[str],
) -> list[Path]:
    """Export slides using the PowerPoint COM automation API (Windows only)."""
    try:
        import win32com.client as _win32  # type: ignore[import-untyped]
    except ImportError as exc:
        raise RuntimeError(
            "pywin32 is required for the PowerPoint backend. "
            "Install it with: pip install pywin32"
        ) from exc

    canvas_w, canvas_h = _canvas_size_from_pptx(pptx_path)
    out_w = canvas_w * scale
    out_h = canvas_h * scale

    pptx_abs = str(pptx_path.resolve())
    pptx_app = None
    prs = None
    generated: list[Path] = []

    try:
        pptx_app = _win32.Dispatch("PowerPoint.Application")
        try:
            pptx_app.Visible = False  # type: ignore[assignment]
        except Exception:
            # PowerPoint is already open and visible — that's fine, just leave it
            pass
        prs = pptx_app.Presentations.Open(pptx_abs, ReadOnly=True, WithWindow=False)  # type: ignore[attr-defined]

        for idx in range(1, prs.Slides.Count + 1):
            slide = prs.Slides(idx)
            title = slide_titles[idx - 1] if idx - 1 < len(slide_titles) else f"Slide {idx}"
            safe = _safe_name(title)
            png_path = out_dir / f"slide-{idx:03d}-{safe}.png"
            slide.Export(str(png_path.resolve()), "PNG", out_w, out_h)
            logger.info("QA pptx [PowerPoint] → %s", png_path.name)
            generated.append(png_path)

    finally:
        if prs is not None:
            try:
                prs.Close()
            except Exception:  # noqa: BLE001
                pass
        if pptx_app is not None:
            try:
                pptx_app.Quit()
            except Exception:  # noqa: BLE001
                pass

    return generated


# ── LibreOffice backend ──────────────────────────────────────────────────────

def _export_libreoffice(
    pptx_path: Path,
    out_dir: Path,
    scale: int,
    slide_titles: list[str],
) -> list[Path]:
    """Export slides using LibreOffice headless --convert-to png."""
    lo_cmd = shutil.which("libreoffice") or shutil.which("soffice")
    if lo_cmd is None:
        raise RuntimeError(
            "LibreOffice is not found on PATH. Install it (>= 7.0) and ensure "
            "'libreoffice' or 'soffice' is accessible."
        )

    canvas_w, canvas_h = _canvas_size_from_pptx(pptx_path)
    # LibreOffice exports at 96 DPI by default; scale is applied by resizing afterwards.
    # We use cairosvg-free PIL/Pillow approach: LO writes PNGs, we resize.
    # If Pillow is unavailable we keep LO output as-is (scale ignored with a warning).

    with tempfile.TemporaryDirectory(prefix="pptx_qa_lo_") as tmp_str:
        tmp_dir = Path(tmp_str)
        cmd = [
            lo_cmd,
            "--headless",
            "--convert-to", "png",
            "--outdir", str(tmp_dir),
            str(pptx_path.resolve()),
        ]
        logger.debug("Running: %s", " ".join(cmd))
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=120,
        )
        if result.returncode != 0:
            raise RuntimeError(
                f"LibreOffice conversion failed (exit {result.returncode}):\n{result.stderr}"
            )

        # LibreOffice names files as "<stem>1.png", "<stem>2.png", …
        stem = pptx_path.stem
        lo_pngs = sorted(
            tmp_dir.glob(f"{stem}*.png"),
            key=lambda p: _lo_slide_index(p.name, stem),
        )

        if not lo_pngs:
            # Fallback: any PNG in tmp_dir
            lo_pngs = sorted(tmp_dir.glob("*.png"), key=lambda p: p.name)

        generated: list[Path] = []
        for idx, lo_png in enumerate(lo_pngs, start=1):
            title = slide_titles[idx - 1] if idx - 1 < len(slide_titles) else f"Slide {idx}"
            safe = _safe_name(title)
            dst = out_dir / f"slide-{idx:03d}-{safe}.png"

            if scale != 1:
                _resize_png(lo_png, dst, canvas_w * scale, canvas_h * scale)
            else:
                shutil.copy2(lo_png, dst)

            logger.info("QA pptx [LibreOffice] → %s", dst.name)
            generated.append(dst)

    return generated


def _lo_slide_index(filename: str, stem: str) -> int:
    """Extract the numeric index from a LibreOffice-generated PNG filename."""
    suffix = filename[len(stem):]
    m = re.search(r"(\d+)", suffix)
    return int(m.group(1)) if m else 0


def _resize_png(src: Path, dst: Path, width: int, height: int) -> None:
    """Resize *src* PNG to (*width* × *height*) and save to *dst*.

    Uses Pillow if available; otherwise copies without resizing (logs a warning).
    """
    try:
        from PIL import Image  # type: ignore[import-untyped]
        with Image.open(src) as img:
            resized = img.resize((width, height), Image.LANCZOS)
            resized.save(str(dst), "PNG")
    except ImportError:
        logger.warning(
            "Pillow not installed — scale=%d ignored for LibreOffice backend. "
            "Install pillow for HiDPI output: pip install pillow",
            width // (width // (width if width else 1)),
        )
        shutil.copy2(src, dst)


# ── Main exporter class ──────────────────────────────────────────────────────

class PptxScreenshotExporter:
    """Render all slides of a ``.pptx`` file to PNG images.

    Args:
        scale:   Output resolution multiplier (default 2 → 2× HiDPI).
        backend: ``"auto"`` | ``"powerpoint"`` | ``"libreoffice"``.
                 ``"auto"`` tries PowerPoint first (Windows), then LibreOffice.
    """

    def __init__(self, scale: int = 2, backend: str = "auto") -> None:
        self.scale = scale
        if backend not in ("auto", "powerpoint", "libreoffice"):
            raise ValueError(f"backend must be 'auto', 'powerpoint', or 'libreoffice'; got {backend!r}")
        self.backend = backend

    def export(self, pptx_path: Path, out_dir: Path) -> list[Path]:
        """Render every slide of *pptx_path* to PNG files in *out_dir*.

        Args:
            pptx_path: Path to the ``.pptx`` source file.
            out_dir:   Destination directory (created if absent).

        Returns:
            List of paths to the generated PNG files, in slide order.

        Raises:
            RuntimeError: If no usable backend is available.
        """
        out_dir.mkdir(parents=True, exist_ok=True)
        slide_titles = _slide_titles_from_pptx(pptx_path)

        resolved_backend = self._resolve_backend()
        logger.debug("PPTX screenshot backend: %s", resolved_backend)

        if resolved_backend == "powerpoint":
            return _export_powerpoint(pptx_path, out_dir, self.scale, slide_titles)
        if resolved_backend == "libreoffice":
            return _export_libreoffice(pptx_path, out_dir, self.scale, slide_titles)

        raise RuntimeError(
            "No PPTX screenshot backend is available.\n"
            "  • Windows: install pywin32 (pip install pywin32) and ensure PowerPoint is installed.\n"
            "  • Linux/macOS: install LibreOffice >= 7.0 and ensure 'libreoffice' is on PATH."
        )

    def _resolve_backend(self) -> str | None:
        """Return the concrete backend to use, or None if none available."""
        if self.backend == "powerpoint":
            return "powerpoint"
        if self.backend == "libreoffice":
            return "libreoffice"
        # auto
        if sys.platform == "win32":
            try:
                import win32com.client  # noqa: F401  # type: ignore[import-untyped]
                return "powerpoint"
            except ImportError:
                pass
        lo = shutil.which("libreoffice") or shutil.which("soffice")
        if lo:
            return "libreoffice"
        return None
