"""
PptxExporter — Export rendered slides to PowerPoint (.pptx).

Consumes the renderer's element tree (list of RenderBox per slide) and writes
shapes to a python-pptx Presentation object.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any
import uuid

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from scripts.rendering.base_card import RenderBox
from scripts.exporting.icon_resolver import IconResolver

logger = logging.getLogger(__name__)

# EMU per pixel at 96 DPI (PowerPoint native)
_EMU_PER_PX = 914400 / 96  # = 9525


def _px(val: float) -> int:
    """Convert a pixel value to EMU (all CSS coordinate tokens are in px)."""
    return int(val * _EMU_PER_PX)


def _rgb(hex_color: str) -> RGBColor | None:
    """Parse ``#RRGGBB`` to ``RGBColor``. Returns None for transparent/invalid."""
    if not hex_color or hex_color.lower() == "transparent":
        return None
    hex_color = hex_color.lstrip("#")
    if len(hex_color) != 6:
        return None
    try:
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16),
        )
    except (ValueError, TypeError):
        return None


# ── Icon font glyph substitution ────────────────────────────────────────────
# Material Icons/Symbols use CSS text ligatures ("campaign" → icon glyph). This
# is a browser-only feature; PowerPoint does NOT process font ligatures, so the
# literal text sits invisible in the slide. We substitute each known ligature name
# with a standard Unicode symbol that renders in any presentation font at export
# time and is always visible in the right colour.
_ICON_UNICODE_MAP: dict[str, str] = {
    # Navigation / directional
    "home":              "⌂",
    "arrow_forward":     "→",
    "arrow_back":        "←",
    "arrow_upward":      "↑",
    "arrow_downward":    "↓",
    "trending_up":       "↑",
    "trending_down":     "↓",
    "open_in_new":       "↗",
    # Status / feedback
    "check":             "✓",
    "check_circle":      "◉",
    "close":             "✕",
    "error":             "✖",
    "warning":           "⚠",
    "info":              "ℹ",
    "help":              "?",
    "flag":              "⚑",
    # People / social
    "person":            "◉",
    "people":            "◎",
    "groups":            "◎",
    "group":             "◎",
    # Actions / tools
    "settings":          "⚙",
    "engineering":       "⚙",
    "build":             "⚙",
    "search":            "⌕",
    "edit":              "✎",
    "delete":            "✕",
    "add":               "+",
    "remove":            "−",
    "campaign":          "◈",
    "rocket_launch":     "▲",
    # Content / media
    "star":              "★",
    "star_border":       "☆",
    "favorite":          "♥",
    "lightbulb":         "◉",
    "label":             "◆",
    "tag":               "◆",
    "bookmark":          "▼",
    # Data / charts
    "bar_chart":         "▤",
    "pie_chart":         "◔",
    "analytics":         "▤",
    "insights":          "▤",
    "table_chart":       "▦",
    "space_dashboard":   "⊞",
    # Geography / globe
    "public":            "◎",
    "language":          "◎",
    "location_on":       "◎",
    "place":             "◎",
    # Industry / transport
    "elevator":          "▲",
    "escalator":         "▲",
    "business":          "◆",
    "factory":           "◆",
    "precision_manufacturing": "⚙",
    # Additional material symbol fallbacks
    "shield":           "🛡",
    "autorenew":        "↻",
    "visibility":       "👁",
    "gavel":            "⚖",
    # Phosphor icon names
    "trend-up":         "↑",
    "trend-down":       "↓",
    "rocket":           "▲",
    "lightbulb":        "◉",
    "chart-bar":        "▤",
    "chart-line":       "▤",
    "chart-scatter":    "▤",
    "tree":             "◆",
    "tree-structure":   "◆",
    "x":                "✕",
    "check":            "✓",
    "intersect":        "◎",
    "gear":             "⚙",
    "magnifying-glass": "⌕",
    "pencil":           "✎",
    "star":             "★",
    "heart":            "♥",
    "warning":          "⚠",
    "info":             "ℹ",
    "question":         "?",
    "flag":             "⚑",
    "house":            "⌂",
    "arrow-right":      "→",
    "arrow-left":       "←",
    "arrow-up":         "↑",
    "arrow-down":       "↓",
    "arrow-up-right":   "↗",
    "plus":             "+",
    "minus":            "−",
    "shield":           "🛡",
    "eye":              "👁",
    "scales":           "⚖",
    "database":         "◆",
    "cpu":              "⚙",
    "cloud":            "◎",
    "lock":             "◆",
    "user":             "◉",
    "users":            "◎",
    "factory":          "◆",
    "globe":            "◎",
    "map-pin":          "◎",
    # Material Symbols — LAPP deck icons
    "account_tree":     "⊕",
    "alt_route":        "⇄",
    "architecture":     "⊟",
    "block":            "⊘",
    "compare_arrows":   "⇔",
    "conversion_path":  "↝",
    "domain":           "⊡",
    "exit_to_app":      "↗",
    "handshake":        "⊗",
    "home_work":        "⌂",
    "hub":              "⊛",
    "inventory_2":      "◫",
    "manufacturing":    "⚙",
    "monitoring":       "◎",
    "quiz":             "?",
    # Material Symbols — EAM deck icons (legacy)
    "smart_toy":        "🤖",
    "psychology":       "🧠",
    "policy":           "📋",
    # Phosphor Icons — EAM deck icons
    "robot":            "🤖",
    "wrench":           "🔧",
    "brain":            "🧠",
    "clipboard-text":   "📋",
    "warning-circle":   "⚠",
    "warning-diamond":  "⚠",
    # Phosphor — common general icons
    "arrow-right":      "→",
    "arrow-left":       "←",
    "check":            "✓",
    "check-circle":     "✓",
    "x":                "✕",
    "x-circle":         "✕",
    "info":             "ℹ",
    "question":         "?",
    "star":             "★",
    "heart":            "♥",
    "user":             "◉",
    "users":            "◎",
    "lock":             "🔒",
    "shield":           "🛡",
    "shield-check":     "🛡",
    "gear":             "⚙",
    "lightning":        "⚡",
    "chart-bar":        "▬",
    "chart-line":       "∿",
    "database":         "⊡",
    "cloud":            "☁",
    "code":             "⟨⟩",
    "file-text":        "📄",
    "folder":           "📁",
    "magnifying-glass": "🔍",
    "envelope":         "✉",
    "phone":            "📞",
    "calendar":         "📅",
    "clock":            "🕐",
    "flag":             "⚑",
    "tag":              "◈",
    "link":             "⊕",
    "globe":            "◎",
    "map-pin":          "◎",
    "buildings":        "⌂",
    "factory":          "◆",
    "cpu":              "⊟",
    "network":          "⊛",
    "rule":             "≡",
    "schema":           "⊹",
    "task_alt":         "✓",
    "topic":            "◆",
    "tune":             "⊟",
    "verified":         "✓",
    "view_timeline":    "☰",
}
# Fallback glyph when the icon name is not in the map above
_ICON_FALLBACK_GLYPH = "◉"
# Font families that use ligatures (substring match, case-insensitive)
_ICON_FONT_MARKERS = ("material icons", "material symbols", "phosphor")

_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}


class PptxExporter:
    """Export rendered slides to a ``.pptx`` file.

    Args:
        project_root: Path to the presentation project folder (for asset resolution).
    """

    def __init__(self, project_root: str | Path | None = None) -> None:
        self.project_root = Path(project_root) if project_root else None
        _cache = (self.project_root / "assets" / "icons") if self.project_root else Path("/tmp/pptx_icon_cache")
        self._icon_resolver = IconResolver(_cache)

    def export(
        self,
        slides: list[RenderBox],
        output_path: str | Path,
        *,
        canvas_width: int = 1280,
        canvas_height: int = 720,
        slide_titles: list[str] | None = None,
        section_breaks: dict[int, str] | None = None,
    ) -> Path:
        """Write the rendered slide boxes to a ``.pptx`` file.

        Args:
            slides: List of ``RenderBox`` objects (one per slide), each containing
                    render elements produced by the layout and card renderers.
            output_path: Destination file path.
            canvas_width: Slide width in px.
            canvas_height: Slide height in px.
            slide_titles: Optional list of slide titles; used as the slide name
                          shown in PowerPoint's slide panel.
            section_breaks: Optional mapping of slide-index → section name;
                            injects PowerPoint sections at those boundaries.

        Returns:
            The resolved output ``Path``.
        """
        prs = Presentation()
        prs.slide_width = _px(canvas_width)
        prs.slide_height = _px(canvas_height)

        blank_layout = prs.slide_layouts[6]  # Blank layout

        for slide_idx, slide_box in enumerate(slides):
            pptx_slide = prs.slides.add_slide(blank_layout)
            # Name the slide after its title so it appears in the slides panel
            title = (
                (slide_titles[slide_idx] if slide_titles and slide_idx < len(slide_titles) else None)
                or f"Slide {slide_idx + 1}"
            )
            pptx_slide._element.cSld.set("name", title)
            for elem in slide_box.elements:
                self._render_element(pptx_slide, elem)

        # Inject PPTX sections if provided
        if section_breaks:
            self._inject_sections(prs, section_breaks)

        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out))
        logger.info("Exported PPTX: %s", out)
        return out

    # ── element rendering ────────────────────────────────────────────────

    @staticmethod
    def _inject_sections(prs: "Presentation", section_breaks: dict[int, str]) -> None:
        """Inject PowerPoint section markers into *prs* using OOXML p14 namespace.

        Sections are a PowerPoint 2010+ feature stored in ``<p14:sectionLst>``
        as a direct child of ``<p:presentation>``.  Each section lists the
        numeric IDs of the slides it contains, taken from ``<p:sldIdLst>``.

        Args:
            prs: The python-pptx ``Presentation`` object (not yet saved).
            section_breaks: Maps slide-index (0-based) → section name.
        """
        from lxml import etree

        P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"
        P14 = f"{{{P14_NS}}}"

        # Collect slide numeric IDs in order from the presentation sldIdLst
        prs_elem = prs._element
        sld_id_lst = prs_elem.find(
            "{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst"
        )
        if sld_id_lst is None:
            logger.warning("Cannot inject PPTX sections: sldIdLst not found")
            return

        slide_ids = [
            int(el.get("id"))
            for el in sld_id_lst
            if el.get("id") is not None
        ]
        n_slides = len(slide_ids)

        # Build ordered list of (start_index, section_name)
        breaks_sorted = sorted(section_breaks.items())

        # Build <p14:sectionLst>
        sec_lst = etree.Element(f"{P14}sectionLst", nsmap={"p14": P14_NS})

        for i, (start_idx, name) in enumerate(breaks_sorted):
            end_idx = breaks_sorted[i + 1][0] if i + 1 < len(breaks_sorted) else n_slides
            sec = etree.SubElement(sec_lst, f"{P14}section")
            sec.set("name", name)
            sec.set("id", "{" + str(uuid.uuid4()).upper() + "}")
            sld_id_lst_sec = etree.SubElement(sec, f"{P14}sldIdLst")
            for idx in range(start_idx, min(end_idx, n_slides)):
                sld_id_el = etree.SubElement(sld_id_lst_sec, f"{P14}sldId")
                sld_id_el.set("id", str(slide_ids[idx]))

        # Append to <p:presentation> — PowerPoint expects it as the last child
        prs_elem.append(sec_lst)
        logger.debug("Injected %d PPTX sections", len(breaks_sorted))

    def _render_element(self, slide, elem: dict[str, Any]) -> None:
        """Dispatch a single render element to the appropriate PPTX shape."""
        etype = elem.get("type", "")
        # Skip shapes with zero or negative width/height — PowerPoint rejects
        # <a:ext cx="0"> and <a:ext cy="0"> as invalid OOXML.
        # Lines use x1/y1/x2/y2 rather than w/h so they are exempt.
        if etype not in ("line", "placeholder"):
            if float(elem.get("w", 1)) <= 0 or float(elem.get("h", 1)) <= 0:
                logger.debug("Skipping zero-dimension element type=%s w=%s h=%s",
                             etype, elem.get("w"), elem.get("h"))
                return
        if etype == "rect":
            self._add_rect(slide, elem)
        elif etype == "text":
            self._add_text(slide, elem)
        elif etype == "line":
            self._add_line(slide, elem)
        elif etype == "image":
            self._add_image(slide, elem)
        elif etype == "icon":
            self._add_icon(slide, elem)
        elif etype == "ellipse":
            self._add_ellipse(slide, elem)
        elif etype == "table":
            self._add_table(slide, elem)
        elif etype == "bullet_list":
            self._add_bullet_list(slide, elem)
        elif etype == "chevron":
            self._add_chevron(slide, elem)
        elif etype == "trapezoid":
            self._add_trapezoid(slide, elem)
        elif etype == "arc":
            self._add_arc(slide, elem)
        elif etype == "placeholder":
            pass  # placeholders are not rendered in PPTX
        else:
            logger.debug("Unknown element type: %s", etype)

    def _add_ellipse(self, slide, elem: dict[str, Any]) -> None:
        """Add an oval auto-shape for elements with ``type == 'ellipse'``."""
        shape = slide.shapes.add_shape(
            9,  # MSO_AUTO_SHAPE_TYPE.OVAL
            _px(elem["x"]),
            _px(elem["y"]),
            _px(elem["w"]),
            _px(elem["h"]),
        )
        fill_color = _rgb(str(elem.get("fill", "#FFFFFF")))
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()
        if hasattr(shape, "shadow"):
            try:
                shape.shadow.inherit = False
                shape.shadow.visible = False
            except Exception:
                pass
        stroke_color = _rgb(str(elem.get("stroke", "")))
        stroke_width = float(elem.get("stroke_width", 1))
        if stroke_color and stroke_width > 0:
            from pptx.util import Pt
            shape.line.color.rgb = stroke_color
            shape.line.width = Pt(stroke_width * 72 / 96)
        else:
            shape.line.fill.background()

    def _add_chevron(self, slide, elem: dict[str, Any]) -> None:
        """Right-pointing chevron (arrow). MSO_SHAPE.CHEVRON = 52."""
        self._add_basic_shape(slide, elem, mso_shape=52)

    def _add_arc(self, slide, elem: dict[str, Any]) -> None:
        """Annular sector (donut slice) using ``MSO_SHAPE.BLOCK_ARC`` (20).

        Required keys: cx, cy, outer_radius, inner_radius, start_angle, end_angle.
        Angles in degrees, 0 = east, CW positive.
        """
        from pptx.oxml.ns import qn
        from lxml import etree as _lE

        cx = float(elem["cx"])
        cy = float(elem["cy"])
        ro = float(elem["outer_radius"])
        ri = float(elem.get("inner_radius", 0))
        start = float(elem["start_angle"])
        end   = float(elem["end_angle"])

        # BLOCK_ARC (20) bounding box = full circle bounding box (2*ro × 2*ro).
        # OOXML blockArc parameters:
        #   adj1 = stAng  — start angle in 60000ths of a degree, 0=east (3 o'clock), CW positive
        #   adj2 = swAng  — swing (arc span) in 60000ths of a degree, positive = CW
        #   adj3 = innerRadius fraction in 50000ths (0=solid pie, 50000=thin outer ring)
        # Our renderer convention: 0=east, clockwise (same as OOXML stAng), so no offset needed.
        adj1 = int((start % 360) * 60000)             # start angle (OOXML east-0, CW)
        adj2 = int(((end - start) % 360) * 60000)     # swing angle (arc span)
        # adj3: PowerPoint's blockArc inner-radius rendering is unreliable across
        # versions.  We always emit solid pie sectors (adj3=0) and rely on the
        # white-ellipse overlay added by each card renderer to create the donut hole.
        adj3 = 0

        shape = slide.shapes.add_shape(
            20,  # MSO_SHAPE.BLOCK_ARC
            _px(cx - ro), _px(cy - ro),
            _px(ro * 2),  _px(ro * 2),
        )

        # Apply adjustment values via OOXML
        try:
            spPr = shape._element.spPr
            prstGeom = spPr.find(qn("a:prstGeom"))
            if prstGeom is not None:
                avLst = prstGeom.find(qn("a:avLst"))
                if avLst is None:
                    avLst = _lE.SubElement(prstGeom, qn("a:avLst"))
                # Remove any existing gd children
                for gd in list(avLst.findall(qn("a:gd"))):
                    avLst.remove(gd)
                for name, val in (("adj1", adj1), ("adj2", adj2), ("adj3", adj3)):
                    gd = _lE.SubElement(avLst, qn("a:gd"))
                    gd.set("name", name)
                    gd.set("fmla", f"val {val}")
        except Exception:
            pass

        fill_color = _rgb(str(elem.get("fill", "#FFFFFF")))
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()
        stroke_color = _rgb(str(elem.get("stroke", "")))
        stroke_width = float(elem.get("stroke_width", 0))
        if stroke_color and stroke_width > 0:
            from pptx.util import Pt
            shape.line.color.rgb = stroke_color
            shape.line.width = Pt(stroke_width * 72 / 96)
        else:
            shape.line.fill.background()

    def _add_trapezoid(self, slide, elem: dict[str, Any]) -> None:
        """Isoceles trapezoid (wider at bottom by default).

        Set ``elem['orientation'] = 'down'`` (default) for a shape that
        narrows toward the top (pyramid layer).
        Set ``orientation = 'up'`` to flip (funnel layer narrowing downward).

        MSO_SHAPE.TRAPEZOID = 3 ("trap"); narrowing is controlled via the OOXML ``adj``
        guide value derived from ``narrow_pct``.
        adj = (1 - narrow_pct) * 50000  — 0 = rectangle, 50000 = triangle.
        """
        orientation = str(elem.get("orientation", "down")).lower()
        shape = self._add_basic_shape(slide, elem, mso_shape=3)
        if orientation == "up":
            shape.rotation = 180.0
        # Apply narrow_pct as OOXML adj on the prstGeom avLst
        narrow_pct = elem.get("narrow_pct")
        if narrow_pct is not None:
            try:
                from pptx.oxml.ns import qn as _qn_t
                from lxml import etree as _lE_t
                adj_val = int((1.0 - float(narrow_pct)) * 50000)
                adj_val = max(0, min(50000, adj_val))
                spPr = shape._element.spPr
                prstGeom = spPr.find(_qn_t("a:prstGeom"))
                if prstGeom is not None:
                    avLst = prstGeom.find(_qn_t("a:avLst"))
                    if avLst is None:
                        avLst = _lE_t.SubElement(prstGeom, _qn_t("a:avLst"))
                    for gd in list(avLst.findall(_qn_t("a:gd"))):
                        avLst.remove(gd)
                    gd = _lE_t.SubElement(avLst, _qn_t("a:gd"))
                    gd.set("name", "adj")
                    gd.set("fmla", f"val {adj_val}")
            except Exception:
                pass

    def _add_basic_shape(self, slide, elem: dict[str, Any], *, mso_shape: int):
        """Shared helper that places an auto-shape and applies fill/stroke
        the same way ``_add_rect`` and ``_add_ellipse`` do."""
        shape = slide.shapes.add_shape(
            mso_shape,
            _px(elem["x"]),
            _px(elem["y"]),
            _px(elem["w"]),
            _px(elem["h"]),
        )
        fill_color = _rgb(str(elem.get("fill", "#FFFFFF")))
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()
        if hasattr(shape, "shadow"):
            try:
                shape.shadow.inherit = False
                shape.shadow.visible = False
            except Exception:
                pass
        stroke_color = _rgb(str(elem.get("stroke", "")))
        stroke_width = float(elem.get("stroke_width", 1))
        if stroke_color and stroke_width > 0:
            from pptx.util import Pt
            shape.line.color.rgb = stroke_color
            shape.line.width = Pt(stroke_width * 72 / 96)
        else:
            shape.line.fill.background()
        return shape

    def _add_rect(self, slide, elem: dict[str, Any]) -> None:
        from pptx.util import Emu

        rx = float(elem.get("rx", 0) or 0)
        w_px = float(elem["w"])
        h_px = float(elem["h"])

        if rx > 0:
            # Use rounded rectangle (MSO_SHAPE 5 = roundRect).
            # OOXML adj = corner_radius / (min(w,h)/2) * 50000 capped at 50000.
            short = min(w_px, h_px)
            adj_val = min(50000, int(rx / (short / 2) * 50000)) if short > 0 else 50000
            shape = slide.shapes.add_shape(
                5,  # MSO_SHAPE.ROUNDED_RECTANGLE
                _px(elem["x"]), _px(elem["y"]),
                _px(w_px), _px(h_px),
            )
            try:
                from pptx.oxml.ns import qn as _rqn
                from lxml import etree as _rE
                spPr = shape._element.spPr
                prstGeom = spPr.find(_rqn("a:prstGeom"))
                if prstGeom is not None:
                    avLst = prstGeom.find(_rqn("a:avLst"))
                    if avLst is None:
                        avLst = _rE.SubElement(prstGeom, _rqn("a:avLst"))
                    for gd in list(avLst.findall(_rqn("a:gd"))):
                        avLst.remove(gd)
                    gd = _rE.SubElement(avLst, _rqn("a:gd"))
                    gd.set("name", "adj")
                    gd.set("fmla", f"val {adj_val}")
            except Exception:
                pass
        else:
            shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE.RECTANGLE
                _px(elem["x"]),
                _px(elem["y"]),
                _px(w_px),
                _px(h_px),
            )
        fill_color = _rgb(str(elem.get("fill", "#FFFFFF")))
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
            # Apply opacity if specified (0-100).
            # OOXML: a:alpha val="100000" = fully opaque, val="0" = fully transparent.
            opacity_pct = elem.get("opacity")
            if opacity_pct is not None:
                try:
                    from pptx.oxml.ns import qn as _qn
                    from lxml import etree as _lET
                    alpha_val = int(float(opacity_pct) * 1000)  # e.g. 40 → 40000
                    spPr = shape._element.spPr
                    solid_fill = spPr.find(".//" + _qn("a:solidFill"))
                    if solid_fill is not None:
                        srgb = solid_fill.find(_qn("a:srgbClr"))
                        if srgb is not None:
                            alpha_el = _lET.SubElement(srgb, _qn("a:alpha"))
                            alpha_el.set("val", str(alpha_val))
                except Exception:
                    pass
        else:
            shape.fill.background()

        if hasattr(shape, "shadow"):
            try:
                shape.shadow.inherit = False
                shape.shadow.visible = False
            except Exception:
                pass

        stroke_color = _rgb(str(elem.get("stroke", "")))
        stroke_width = float(elem.get("stroke_width", 1))
        if stroke_color and stroke_width > 0:
            shape.line.color.rgb = stroke_color
            shape.line.width = Pt(stroke_width * 72 / 96)
        else:
            shape.line.fill.background()

    def _add_text(self, slide, elem: dict[str, Any]) -> None:
        from pptx.util import Emu
        w = max(1.0, float(elem.get("w", 200)))
        h = max(1.0, float(elem.get("h", elem.get("font_size", 14) * 2)))
        txBox = slide.shapes.add_textbox(
            _px(elem["x"]),
            _px(elem["y"]),
            _px(w),
            _px(h),
        )
        try:
            txBox.line.fill.background()
            if hasattr(txBox, "shadow"):
                txBox.shadow.inherit = False
                txBox.shadow.visible = False
        except Exception:
            pass
        tf = txBox.text_frame
        # Zero internal margins so element coordinates match visual positions exactly
        tf.margin_top = 0
        tf.margin_bottom = 0
        tf.margin_left = 0
        tf.margin_right = 0
        tf.word_wrap = elem.get("wrap", False)

        # Vertical alignment via XML bodyPr anchor attribute
        v_align = elem.get("vertical_align", "top")
        _ANCHOR_MAP = {"top": "t", "middle": "ctr", "bottom": "b"}
        anchor_val = _ANCHOR_MAP.get(v_align, "t")
        _NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
        body_pr = tf._txBody.find(f"{{{_NS}}}bodyPr")
        if body_pr is not None:
            body_pr.set("anchor", anchor_val)

        # ── Icon font substitution ────────────────────────────────────────
        # Material Symbols/Icons uses CSS text ligatures which PPTX does not
        # process. Substitute the ligature name with a standard Unicode symbol
        # and drop the special font requirement so the glyph is always visible.
        raw_text = str(elem.get("text", ""))
        font_family = elem.get("font_family")
        if font_family and any(m in font_family.lower() for m in _ICON_FONT_MARKERS):
            # Ligature name heuristic: all-lowercase, may contain underscores or hyphens
            if raw_text and (raw_text.islower() or "_" in raw_text or "-" in raw_text):
                raw_text = _ICON_UNICODE_MAP.get(raw_text, _ICON_FALLBACK_GLYPH)
                font_family = None  # standard font — glyph is in BMP, no icon font needed

        p = tf.paragraphs[0]
        p.alignment = _ALIGN_MAP.get(elem.get("alignment", "left"), PP_ALIGN.LEFT)
        line_height = elem.get("line_height")
        if line_height is not None:
            try:
                p.line_spacing = Pt(float(line_height))
            except Exception:
                pass

        # Font-size tokens are defined in pt — use directly (no px→pt conversion needed)
        font_size_pt = Pt(float(elem.get("font_size", 14)))
        color = _rgb(str(elem.get("font_color", "#000000")))
        weight = elem.get("font_weight", "normal")
        style = str(elem.get("font_style") or "").lower()
        is_bold = str(weight).lower() in ("bold", "700") or "bold" in style
        is_italic = "italic" in style

        runs_data = elem.get("runs")
        if runs_data:
            # Multi-run path — inline bold/italic markup was present.
            # Do NOT use p.text setter; add each run explicitly so per-run
            # bold/italic attributes are preserved.
            for r_data in runs_data:
                run = p.add_run()
                run.text = r_data["text"]
                run.font.size = font_size_pt
                if color:
                    run.font.color.rgb = color
                run.font.bold = is_bold or bool(r_data.get("bold"))
                run.font.italic = is_italic or bool(r_data.get("italic"))
                if font_family:
                    run.font.name = str(font_family)
        else:
            # Single-run path — plain text, no inline markup.
            p.text = raw_text
            # Apply font properties to ALL runs — p.text setter splits on \n and
            # creates multiple <a:r> runs via add_br(); only setting run[0] would
            # leave subsequent runs with the slide-default (large) font size.
            for r in p.runs:
                r.font.size = font_size_pt
                if color:
                    r.font.color.rgb = color
                r.font.bold = is_bold
                r.font.italic = is_italic
                if font_family:
                    r.font.name = str(font_family)

    def _add_bullet_list(self, slide, elem: dict[str, Any]) -> None:
        """Render a bullet list as a single PPTX text box with native bullet paragraph formatting.

        Each bullet item becomes a paragraph with hanging-indent layout and
        ``<a:buChar>`` / ``<a:buClr>`` OOXML attributes so PowerPoint renders
        the marker natively — no separate marker text boxes needed.
        """
        from lxml import etree
        from pptx.oxml.ns import qn

        txBox = slide.shapes.add_textbox(
            _px(elem["x"]), _px(elem["y"]),
            _px(elem["w"]), _px(elem["h"]),
        )
        try:
            txBox.line.fill.background()
            if hasattr(txBox, "shadow"):
                txBox.shadow.inherit = False
                txBox.shadow.visible = False
        except Exception:
            pass

        tf = txBox.text_frame
        tf.margin_top = 0
        tf.margin_bottom = 0
        tf.margin_left = 0
        tf.margin_right = 0
        tf.word_wrap = True

        font_size_pt = Pt(float(elem.get("font_size", 14)))
        font_color_rgb = _rgb(str(elem.get("font_color") or "#000000"))
        weight = elem.get("font_weight", "normal")
        style_str = str(elem.get("font_style") or "").lower()
        is_bold = str(weight).lower() in ("bold", "700") or "bold" in style_str
        is_italic = "italic" in style_str
        align = _ALIGN_MAP.get(elem.get("alignment", "left"), PP_ALIGN.LEFT)

        bullet_char = str(elem.get("bullet_char") or "\u2022")  # fallback: •
        bullet_color_hex = str(
            elem.get("bullet_color") or elem.get("font_color") or "#000000"
        ).lstrip("#")
        bullet_indent_px = float(elem.get("bullet_indent", 16))
        bullet_gap_px = float(elem.get("bullet_gap", 8))
        bullet_spacing_val = float(elem.get("bullet_spacing", 4))
        # Total hanging indent = marker column + gap between marker and text
        total_hanging_emu = int((bullet_indent_px + bullet_gap_px) * _EMU_PER_PX)
        # spcAft in hundredths of a point (1 px ≈ 0.75 pt → × 75)
        spc_aft_hpt = int(bullet_spacing_val * 75)

        items = elem.get("items", [])
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align

            # Hanging-indent layout: marL pushes text right, indent pulls marker left
            pPr = p._p.get_or_add_pPr()
            pPr.set("marL", str(total_hanging_emu))
            pPr.set("indent", str(-total_hanging_emu))

            # Vertical spacing below each item
            if spc_aft_hpt > 0:
                spcAft = etree.SubElement(pPr, qn("a:spcAft"))
                spcPts_el = etree.SubElement(spcAft, qn("a:spcPts"))
                spcPts_el.set("val", str(spc_aft_hpt))

            # Bullet colour  <a:buClr><a:srgbClr val="RRGGBB"/></a:buClr>
            if len(bullet_color_hex) == 6:
                buClr = etree.SubElement(pPr, qn("a:buClr"))
                srgbClr = etree.SubElement(buClr, qn("a:srgbClr"))
                srgbClr.set("val", bullet_color_hex)

            # Bullet character (or suppress with buNone when style=none)
            if bullet_char:
                buChar = etree.SubElement(pPr, qn("a:buChar"))
                buChar.set("char", bullet_char)
            else:
                etree.SubElement(pPr, qn("a:buNone"))

            # Text run(s)
            runs_data = item.get("runs")
            if runs_data:
                for r_data in runs_data:
                    run = p.add_run()
                    run.text = r_data["text"]
                    run.font.size = font_size_pt
                    if font_color_rgb:
                        run.font.color.rgb = font_color_rgb
                    run.font.bold = is_bold or bool(r_data.get("bold"))
                    run.font.italic = is_italic or bool(r_data.get("italic"))
            else:
                run = p.add_run()
                run.text = str(item.get("text", ""))
                run.font.size = font_size_pt
                if font_color_rgb:
                    run.font.color.rgb = font_color_rgb
                run.font.bold = is_bold
                run.font.italic = is_italic

    def _add_line(self, slide, elem: dict[str, Any]) -> None:
        connector = slide.shapes.add_connector(
            1,  # MSO_CONNECTOR.STRAIGHT
            _px(elem["x1"]),
            _px(elem["y1"]),
            _px(elem["x2"]),
            _px(elem["y2"]),
        )
        stroke_color = _rgb(str(elem.get("stroke", "")))
        if stroke_color:
            connector.line.color.rgb = stroke_color
            connector.line.width = Pt(float(elem.get("stroke_width", 1)) * 72 / 96)
            if elem.get("dashed"):
                try:
                    from pptx.oxml.ns import qn as _qn2
                    from lxml import etree as _et2
                    spPr = connector._element.spPr
                    ln_el = spPr.find(_qn2("a:ln"))
                    if ln_el is None:
                        ln_el = _et2.SubElement(spPr, _qn2("a:ln"))
                    pd = ln_el.find(_qn2("a:prstDash"))
                    if pd is None:
                        pd = _et2.SubElement(ln_el, _qn2("a:prstDash"))
                    pd.set("val", "dash")
                except Exception:
                    pass
        else:
            connector.line.fill.background()
        try:
            if hasattr(connector, "shadow"):
                connector.shadow.inherit = False
                connector.shadow.visible = False
        except Exception:
            pass

    def _add_svg_pic(
        self, slide, svg_path: "Path", x: float, y: float, w: float, h: float
    ) -> bool:
        """Embed an SVG as a native Office 365 picture shape in PPTX.

        Bypasses PIL/cairo entirely by injecting the SVG bytes directly as an
        OPC image part with ``image/svg+xml`` content type and wiring a
        ``<p:pic>`` OOXML element into the slide's shape tree.

        Office 365 / PowerPoint 2019+ renders SVG natively.  Older versions
        will display a blank, but this is acceptable because they also lacked
        native SVG support regardless of the embedding approach.

        Returns *True* on success, *False* on any failure.
        """
        try:
            from lxml import etree as _lE
            from pptx.opc.constants import RELATIONSHIP_TYPE as RT
            from pptx.parts.image import ImagePart

            svg_bytes = svg_path.read_bytes()
            slide_part = slide.part
            package = slide_part._package

            # Allocate a new /ppt/media/imageN.svg part
            partname = package.next_image_partname("svg")
            img_part = ImagePart(partname, "image/svg+xml", package, svg_bytes)
            rId = slide_part.relate_to(img_part, RT.IMAGE)

            # Unique shape id (max of existing ids + 1)
            sp_tree = slide.shapes._spTree
            existing_ids = {
                int(el.get("id", 0))
                for el in sp_tree.iter()
                if el.get("id") is not None
            }
            shape_id = max(existing_ids, default=0) + 1

            emu_x = int(_px(x))
            emu_y = int(_px(y))
            emu_w = int(_px(w))
            emu_h = int(_px(h))

            PNS = "http://schemas.openxmlformats.org/presentationml/2006/main"
            ANS = "http://schemas.openxmlformats.org/drawingml/2006/main"
            RNS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

            pic_xml = (
                f'<p:pic xmlns:p="{PNS}" xmlns:a="{ANS}" xmlns:r="{RNS}">'
                f'<p:nvPicPr>'
                f'<p:cNvPr id="{shape_id}" name="Diagram {shape_id}"/>'
                f'<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>'
                f'<p:nvPr/>'
                f'</p:nvPicPr>'
                f'<p:blipFill>'
                f'<a:blip r:embed="{rId}"/>'
                f'<a:stretch><a:fillRect/></a:stretch>'
                f'</p:blipFill>'
                f'<p:spPr>'
                f'<a:xfrm><a:off x="{emu_x}" y="{emu_y}"/>'
                f'<a:ext cx="{emu_w}" cy="{emu_h}"/></a:xfrm>'
                f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
                f'</p:spPr>'
                f'</p:pic>'
            )
            sp_tree.append(_lE.fromstring(pic_xml.encode()))
            logger.debug("Native SVG embedded: %s rId=%s", svg_path.name, rId)
            return True

        except Exception as exc:
            logger.warning("SVG pic failed for %s: %s", getattr(svg_path, 'name', str(svg_path)), exc)
            return False

    def _add_icon(self, slide, elem: dict[str, Any]) -> None:
        """Render an icon element — direct SVG picture (Office 365) or Unicode fallback."""
        name = str(elem.get("name", ""))
        color = str(elem.get("color") or "#000000")
        font_family = str(elem.get("font_family") or "")
        x, y = elem["x"], elem["y"]
        size = float(elem.get("w") or elem.get("h") or 20)
        font_size = float(elem.get("font_size") or size)

        icon_family_lower = font_family.lower()
        is_icon_font = any(
            marker in icon_family_lower
            for marker in ("material icons", "material symbols", "phosphor")
        )

        # Strategy 1: direct SVG picture (Office 365 native SVG support)
        if is_icon_font:
            svg_path = self._icon_resolver.resolve(name, font_family, color)
            if svg_path is not None and self._add_svg_pic(slide, svg_path, x, y, size, size):
                return

        # Strategy 2: Unicode fallback
        fallback_char = _ICON_UNICODE_MAP.get(name, _ICON_FALLBACK_GLYPH)
        self._add_text(slide, {
            "type": "text",
            "x": x, "y": y, "w": size, "h": size,
            "text": fallback_char,
            "font_size": font_size,
            "font_color": color,
            "font_weight": "bold",
            "alignment": "center",
            "vertical_align": "middle",
        })

    @staticmethod
    def _svg_viewbox_dims(img_path: str) -> tuple[float, float] | tuple[None, None]:
        """Parse an SVG file's natural dimensions from its viewBox or width/height attrs.

        Returns (width, height) in user units, or (None, None) if not determinable.
        """
        try:
            import xml.etree.ElementTree as _ET
            root = _ET.parse(img_path).getroot()
            # Strip namespace prefix from tag if present
            ns_prefix = "{http://www.w3.org/2000/svg}"
            # Try viewBox first (most reliable)
            vb = root.get("viewBox", "")
            if vb:
                parts = vb.replace(",", " ").split()
                if len(parts) >= 4:
                    return float(parts[2]), float(parts[3])
            # Fall back to width/height attributes
            w_attr = root.get("width", "")
            h_attr = root.get("height", "")
            if w_attr and h_attr:
                return float(w_attr.rstrip("px ")), float(h_attr.rstrip("px "))
        except Exception:
            pass
        return None, None

    def _add_image(self, slide, elem: dict[str, Any]) -> None:
        src = elem.get("src", "")
        if not src:
            # Render placeholder rect instead
            self._add_rect(
                slide,
                {
                    "x": elem["x"],
                    "y": elem["y"],
                    "w": elem["w"],
                    "h": elem["h"],
                    "fill": "#F5F5F5",
                    "stroke": "#CCCCCC",
                    "stroke_width": 1,
                },
            )
            return

        # Resolve path
        img_path = None
        if self.project_root:
            candidate = self.project_root / "assets" / src
            if candidate.exists():
                img_path = str(candidate)

        if img_path:
            try:
                picture = None

                target_w = float(elem["w"])
                target_h = float(elem["h"])
                fit_x = float(elem["x"])
                fit_y = float(elem["y"])
                fit_w = target_w
                fit_h = target_h
                fit_mode = elem.get("fit", "contain")

                if img_path.lower().endswith(".svg"):
                    # SVG: read natural dims from viewBox; let the SVG's own
                    # preserveAspectRatio handle the internal rendering.
                    # For cover mode: place at full target box — the SVG's
                    #   preserveAspectRatio="xMidYMid slice" fills the box,
                    #   clipping edges symmetrically. No letterboxing.
                    # For contain mode: scale so the SVG fits within the target
                    #   box while maintaining ratio; center with letterboxing.
                    if fit_mode == "contain":
                        nat_w, nat_h = self._svg_viewbox_dims(img_path)
                        if nat_w and nat_h and nat_w > 0 and nat_h > 0:
                            scale = min(target_w / nat_w, target_h / nat_h)
                            fit_w = nat_w * scale
                            fit_h = nat_h * scale
                            fit_x = float(elem["x"]) + (target_w - fit_w) / 2
                            fit_y = float(elem["y"]) + (target_h - fit_h) / 2
                            fit_x = float(elem["x"]) + (target_w - fit_w) / 2
                            fit_y = float(elem["y"]) + (target_h - fit_h) / 2
                    # cover: fit_w/fit_h stay at target_w/target_h (full box)

                    from pathlib import Path as _Path
                    svg_ok = self._add_svg_pic(slide, _Path(img_path), fit_x, fit_y, fit_w, fit_h)
                    if not svg_ok:
                        # Fallback: labeled placeholder so the layout is not broken
                        self._add_rect(slide, {
                            "x": fit_x, "y": fit_y, "w": fit_w, "h": fit_h,
                            "fill": "#F5F5F5", "stroke": "#CCCCCC", "stroke_width": 1,
                        })
                        self._add_text(slide, {
                            "x": fit_x, "y": fit_y + fit_h * 0.4,
                            "w": fit_w, "h": fit_h * 0.2,
                            "text": f"[{_Path(img_path).stem}]",
                            "font_size": 11,
                            "font_color": "#999999",
                            "alignment": "center",
                        })
                    picture = None  # _add_svg_pic/_fallback handle their own output
                else:
                    # Raster: use PIL for natural dims + apply fit mode
                    nat_w, nat_h = None, None
                    try:
                        from PIL import Image as _PILImg
                        with _PILImg.open(img_path) as _img:
                            nat_w, nat_h = _img.size
                    except Exception:
                        pass

                    if nat_w and nat_h and nat_w > 0 and nat_h > 0:
                        if fit_mode == "cover":
                            scale = max(target_w / nat_w, target_h / nat_h)
                        else:  # contain (default)
                            scale = min(target_w / nat_w, target_h / nat_h)
                        fit_w = nat_w * scale
                        fit_h = nat_h * scale
                        fit_x = float(elem["x"]) + (target_w - fit_w) / 2
                        fit_y = float(elem["y"]) + (target_h - fit_h) / 2

                    picture = slide.shapes.add_picture(
                        img_path,
                        _px(fit_x),
                        _px(fit_y),
                        _px(fit_w),
                        _px(fit_h),
                    )
                if picture is not None:
                    try:
                        picture.line.fill.background()
                        if hasattr(picture, "shadow"):
                            picture.shadow.inherit = False
                            picture.shadow.visible = False
                    except Exception:
                        pass
            except Exception as exc:
                logger.warning("Failed to add image %s to PPTX: %s", src, exc)
        else:
            logger.warning("Image not found: %s — rendering placeholder", src)
            # Placeholder with missing path text
            self._add_rect(
                slide,
                {
                    "x": elem["x"],
                    "y": elem["y"],
                    "w": elem["w"],
                    "h": elem["h"],
                    "fill": "#FFF3E0",
                    "stroke": "#FF9800",
                    "stroke_width": 1,
                },
            )
            self._add_text(
                slide,
                {
                    "x": elem["x"],
                    "y": elem["y"] + elem["h"] * 0.4,
                    "w": elem["w"],
                    "text": f"[missing: {src}]",
                    "font_size": 10,
                    "font_color": "#FF9800",
                    "alignment": "center",
                },
            )

    def _add_table(self, slide, elem: dict[str, Any]) -> None:
        """Render a ``table`` element as a native python-pptx table shape.

        The table is positioned at the coordinates given in *elem* and sized
        to fit the available body area.  Each row descriptor in ``elem['rows']``
        carries its own styling so header, data, stripe, and sum rows are all
        handled uniformly.
        """
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        all_rows: list[dict] = elem.get("rows", [])
        col_widths_px: list[float] = elem.get("col_widths", [])
        border_color_hex: str = str(elem.get("border_color") or "#E5E7EB")
        border_width_pt: float = float(elem.get("border_width") or 1) * 72 / 96
        pad_x: float = float(elem.get("pad_x") or 8)
        pad_y: float = float(elem.get("pad_y") or 4)

        n_rows = len(all_rows)
        n_cols = len(col_widths_px)
        if n_rows == 0 or n_cols == 0:
            return

        # Total height: sum of row heights (after any scaling applied in table_card)
        total_height = sum(float(r.get("row_height", 24)) for r in all_rows)
        available_h = float(elem.get("h") or total_height)
        # Scale row heights proportionally only when they exceed available space.
        # Use total_height (not available_h) for the shape so PowerPoint doesn't
        # stretch rows to fill a larger container than the content needs.
        scale = min(1.0, available_h / total_height) if total_height > 0 else 1.0

        row_heights_px = [float(r.get("row_height", 24)) * scale for r in all_rows]
        shape_height = sum(row_heights_px)  # actual content height, not box.h

        # python-pptx wants col widths and row heights in EMU
        col_w_emu = [_px(w) for w in col_widths_px]
        row_h_emu = [_px(h) for h in row_heights_px]

        tbl = slide.shapes.add_table(
            n_rows,
            n_cols,
            _px(elem["x"]),
            _px(elem["y"]),
            _px(elem["w"]),
            _px(shape_height),  # shape sized to content, not full card slot
        ).table

        # Set column widths
        for ci, cw in enumerate(col_w_emu):
            tbl.columns[ci].width = cw

        # Set row heights
        for ri, rh in enumerate(row_h_emu):
            tbl.rows[ri].height = rh

        # Style each cell
        for ri, row_desc in enumerate(all_rows):
            cells: list[str] = row_desc.get("cells", [])
            bg_hex: str = str(row_desc.get("bg_color") or "#FFFFFF")
            fg_hex: str = str(row_desc.get("font_color") or "#000000")
            font_size: float = float(row_desc.get("font_size") or 12)
            weight: str = str(row_desc.get("font_weight") or "normal")
            style: str = str(row_desc.get("font_style") or "normal").lower()
            aligns: list[str] = row_desc.get("alignments") or []
            border_bottom_hex: str = str(row_desc.get("border_bottom_color") or border_color_hex)

            is_bold = str(weight).lower() in ("bold", "700") or "bold" in style
            is_italic = "italic" in style

            for ci in range(n_cols):
                cell = tbl.cell(ri, ci)
                cell_text = cells[ci] if ci < len(cells) else ""
                alignment = aligns[ci] if ci < len(aligns) else "left"

                # Background fill
                bg_rgb = _rgb(bg_hex)
                if bg_rgb:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = bg_rgb
                else:
                    cell.fill.background()

                # Text frame
                tf = cell.text_frame
                tf.margin_top = _px(pad_y)
                tf.margin_bottom = _px(pad_y)
                tf.margin_left = _px(pad_x)
                tf.margin_right = _px(pad_x)
                tf.word_wrap = True

                para = tf.paragraphs[0]
                para.text = cell_text
                para.alignment = _ALIGN_MAP.get(alignment, PP_ALIGN.LEFT)

                if para.runs:
                    run = para.runs[0]
                else:
                    run = para.add_run()
                    run.text = cell_text

                run.font.size = Pt(font_size)
                run.font.bold = is_bold
                run.font.italic = is_italic
                fg_rgb = _rgb(fg_hex)
                if fg_rgb:
                    run.font.color.rgb = fg_rgb

                # Cell borders via tcPr XML — apply uniform border_color
                self._set_table_cell_borders(
                    cell, border_color_hex, border_width_pt,
                    border_bottom_hex if ri == 0 else border_color_hex,
                )

    @staticmethod
    def _set_table_cell_borders(
        cell: Any,
        border_color_hex: str,
        border_width_pt: float,
        bottom_color_hex: str,
    ) -> None:
        """Apply solid borders on all four sides of a python-pptx table cell via OOXML."""
        from lxml import etree
        from pptx.oxml.ns import qn

        def _make_ln(color_hex: str, width_pt: float) -> "etree._Element":
            w_emu = int(width_pt * 12700)  # pt → EMU (1 pt = 12700 EMU)
            ln = etree.Element(qn("a:ln"), attrib={"w": str(w_emu)})
            solid = etree.SubElement(ln, qn("a:solidFill"))
            srgb = etree.Element(qn("a:srgbClr"), attrib={"val": color_hex.lstrip("#")})
            solid.append(srgb)
            return ln

        tc = cell._tc
        tcPr = tc.find(qn("a:tcPr"))
        if tcPr is None:
            tcPr = etree.SubElement(tc, qn("a:tcPr"))

        for tag, color in (
            ("a:lnL", border_color_hex),
            ("a:lnR", border_color_hex),
            ("a:lnT", border_color_hex),
            ("a:lnB", bottom_color_hex),
        ):
            existing = tcPr.find(qn(tag))
            if existing is not None:
                tcPr.remove(existing)
            tcPr.append(_make_ln(color, border_width_pt))
