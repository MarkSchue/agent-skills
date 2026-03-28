"""
visual_extractor.py — Derive design tokens from images, PPTX, or URLs
─────────────────────────────────────────────────────────────────────────────
Analyzes a visual source and proposes semantic design token mappings.
Before writing anything, results are presented to the user for approval.

Supported sources:
  - PNG / JPG / WEBP image files
  - .pptx files (extracts theme colors + fonts)
  - .drawio files (extracts style fill/stroke colors + fonts)
  - URLs (fetches CSS computed styles; live URL)

Usage:
    python scripts/visual_extractor.py <source> [--output design-config.yaml]

Requires: pyyaml, Pillow (images), python-pptx (pptx), requests (URLs)
"""

from __future__ import annotations
import argparse
import sys
from pathlib import Path
from typing import Any


# ── Extraction strategies ─────────────────────────────────────────────────────

def extract_from_image(path: Path) -> dict[str, Any]:
    """Extract dominant colors from a PNG/JPG using Pillow."""
    try:
        from PIL import Image
    except ImportError:
        print("ERROR: Pillow not installed. Run: pip install Pillow", file=sys.stderr)
        return {}

    img = Image.open(path).convert("RGB")
    img_small = img.resize((200, 200))  # Downsample for speed

    # Count color frequencies
    pixels = [(r, g, b) for r, g, b in img_small.getdata()]  # type: ignore[misc]
    from collections import Counter
    color_counts = Counter(pixels)

    # Group into buckets (round to nearest 16 to reduce noise)
    bucketed: Counter = Counter()
    for (r, g, b), count in color_counts.items():
        key = (r // 16 * 16, g // 16 * 16, b // 16 * 16)
        bucketed[key] += count

    # Top 10 colors
    top_colors = bucketed.most_common(15)

    # Map to semantic roles heuristically
    result: dict[str, str] = {}
    _assign_semantic_roles(result, top_colors)
    return {"colors": result, "source": str(path)}


def extract_from_pptx(path: Path) -> dict[str, Any]:
    """Extract color theme and fonts from a .pptx file."""
    try:
        import pptx
    except ImportError:
        print("ERROR: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
        return {}

    from pptx import Presentation
    prs = Presentation(str(path))

    colors: dict[str, str] = {}
    fonts: list[str] = []

    # Extract theme colors
    theme = prs.core_properties
    slide_width  = prs.slide_width
    slide_height = prs.slide_height

    # Scan first 3 slides for fill/text colors
    for slide in list(prs.slides)[:3]:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:  # type: ignore[union-attr]
                    for run in para.runs:
                        if run.font.color and run.font.color.type:
                            try:
                                rgb = run.font.color.rgb
                                hex_color = f"#{rgb:06x}"
                                if hex_color not in fonts:
                                    fonts.append(hex_color)
                            except Exception:
                                pass
                        if run.font.name:
                            if run.font.name not in fonts:
                                fonts.append(run.font.name)
            if hasattr(shape, "fill") and shape.fill.type is not None:  # type: ignore[union-attr]
                try:
                    from pptx.enum.dml import MSO_THEME_COLOR
                    rgb = shape.fill.fore_color.rgb  # type: ignore[union-attr]
                    hex_color = f"#{rgb:06x}"
                    if hex_color not in colors.values():
                        colors[f"color_{len(colors)}"] = hex_color
                except Exception:
                    pass

    return {
        "colors":   colors,
        "fonts":    fonts[:5],
        "canvas":   {"width": (slide_width or 0) // 9525, "height": (slide_height or 0) // 9525},
        "source":   str(path),
    }


def extract_from_drawio(path: Path) -> dict[str, Any]:
    """Extract fill/stroke colors and fonts from a .drawio XML file."""
    import re
    from xml.etree import ElementTree as ET

    try:
        tree = ET.parse(str(path))
    except ET.ParseError as e:
        print(f"ERROR parsing draw.io file: {e}", file=sys.stderr)
        return {}

    hex_pattern = re.compile(r"(?:fillColor|strokeColor|fontColor)=(#[0-9a-fA-F]{6})")
    font_pattern = re.compile(r"fontFamily=([^;\"]+)")

    colors: list[str] = []
    fonts:  list[str] = []

    for elem in tree.iter():
        style = elem.get("style", "")
        colors.extend(hex_pattern.findall(style))
        fonts.extend(font_pattern.findall(style))

    from collections import Counter
    top_colors = [c for c, _ in Counter(colors).most_common(12)]
    top_fonts  = list(dict.fromkeys(fonts))[:3]

    result: dict[str, str] = {}
    _assign_semantic_roles_from_list(result, top_colors)
    return {"colors": result, "fonts": top_fonts, "source": str(path)}


def extract_from_url(url: str) -> dict[str, Any]:
    """Fetch a URL and extract CSS color values from inline styles and stylesheets."""
    try:
        import requests
    except ImportError:
        print("ERROR: requests not installed. Run: pip install requests", file=sys.stderr)
        return {}

    import re
    resp = requests.get(url, timeout=15, headers={"User-Agent": "Mozilla/5.0"})
    if not resp.ok:
        print(f"ERROR: HTTP {resp.status_code} fetching {url}", file=sys.stderr)
        return {}

    hex_pattern = re.compile(r"#(?:[0-9a-fA-F]{3}|[0-9a-fA-F]{6})\b")
    colors_found = hex_pattern.findall(resp.text)

    from collections import Counter
    top_colors = [c.lower() for c, _ in Counter(colors_found).most_common(20)
                  if c.lower() not in ("#fff", "#ffffff", "#000", "#000000")]

    result: dict[str, str] = {}
    _assign_semantic_roles_from_list(result, top_colors[:10])
    return {"colors": result, "source": url}


# ── Semantic role assignment heuristics ───────────────────────────────────────

def _assign_semantic_roles(result: dict, top_colors: list) -> None:
    """Assign roles based on brightness/position in the color list."""
    roles = ["surface", "on-surface", "primary", "secondary", "accent",
             "neutral", "error", "warning", "success", "on-primary"]
    for i, ((r, g, b), _count) in enumerate(top_colors[:len(roles)]):
        hex_color = f"#{r:02x}{g:02x}{b:02x}"
        result[roles[i]] = hex_color


def _assign_semantic_roles_from_list(result: dict, colors: list[str]) -> None:
    roles = ["primary", "secondary", "accent", "neutral", "surface",
             "on-surface", "error", "warning", "success", "on-primary"]
    for i, color in enumerate(colors[:len(roles)]):
        result[roles[i]] = color


# ── Presentation / approval workflow ─────────────────────────────────────────

def present_findings(findings: dict) -> None:
    """Print extraction findings in a format suitable for user review."""
    print("\n=== Visual Extraction Results ===")
    print(f"Source: {findings.get('source', 'unknown')}\n")

    colors = findings.get("colors", {})
    if colors:
        print("Proposed color token mappings:")
        print("(Review each — rename, reject, or accept)\n")
        for token, hex_val in colors.items():
            print(f"  {token:15s} → {hex_val}")
    else:
        print("  No colors extracted.")

    fonts = findings.get("fonts", [])
    if fonts:
        print(f"\nDetected fonts: {', '.join(str(f) for f in fonts)}")

    canvas = findings.get("canvas")
    if canvas:
        print(f"\nCanvas: {canvas.get('width')} × {canvas.get('height')} px")

    print("\n⚠ No files will be written until you approve these mappings.")
    print("Run with  --write  to write approved tokens to design-config.yaml\n")


# ── Main ──────────────────────────────────────────────────────────────────────

def main(source: str, output_path: Path | None, write: bool) -> None:
    path = Path(source)

    if source.startswith("http://") or source.startswith("https://"):
        findings = extract_from_url(source)
    elif path.suffix.lower() in (".png", ".jpg", ".jpeg", ".webp"):
        findings = extract_from_image(path)
    elif path.suffix.lower() == ".pptx":
        findings = extract_from_pptx(path)
    elif path.suffix.lower() == ".drawio":
        findings = extract_from_drawio(path)
    else:
        print(f"ERROR: Unsupported source type: {path.suffix}", file=sys.stderr)
        sys.exit(1)

    present_findings(findings)

    if write and output_path and findings.get("colors"):
        import yaml
        config = {}
        if output_path.exists():
            with open(output_path) as fh:
                config = yaml.safe_load(fh) or {}
        config.setdefault("colors", {}).update(findings["colors"])
        if findings.get("canvas"):
            config.setdefault("canvas", {}).update(findings["canvas"])
        with open(output_path, "w") as fh:
            yaml.dump(config, fh, default_flow_style=False, allow_unicode=True)
        print(f"Wrote approved tokens to {output_path}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Extract design tokens from a visual source")
    parser.add_argument("source",   help="PNG/JPG/PPTX/drawio path or https:// URL")
    parser.add_argument("--output", help="design-config.yaml to write tokens into", default=None)
    parser.add_argument("--write",  action="store_true",
                        help="Write approved tokens (skip for dry-run preview)")
    args = parser.parse_args()
    main(args.source, Path(args.output) if args.output else None, args.write)
