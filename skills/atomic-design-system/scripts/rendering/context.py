"""
rendering/context.py — Unified render context for draw.io and PPTX backends
============================================================================
Both DrawioCtx and PptxCtx expose an identical public API so every molecule
class is written ONCE and works for both output formats.

Public API shared by both backends
-----------------------------------
    ctx.rect(x, y, w, h, fill, stroke=None, radius=0)
    ctx.text(x, y, w, h, text, size=14, bold=False, italic=False,
             color=None, align="left", valign="top", inner_margin=0,
             font=None)
    ctx.badge(x, y, w, h, text, fill, text_color=None, size=11, radius=12)
    ctx.divider(x, y, w, color=None)
    ctx.ellipse(x, y, w, h, fill, stroke=None)
    ctx.line(x1, y1, x2, y2, color=None)

    ctx.PAD          → int  (canvas padding from CSS --canvas-padding)
    ctx.font         → str  (font-family from CSS --font-family)
    ctx.color(token) → str  (hex color from CSS --color-{token})
    ctx.rad(key)     → int  (border radius from CSS --{key})
    ctx.status_color(status)  → (bg_hex, fg_hex)
    ctx.progress_color(pct)   → hex color
    ctx.card_bg_color(props, default_token) → hex  (semantic card bg with deck.md override via card_bg: filled|clean|alt|featured)

Per-Card Instance Overrides — Canonical Pattern
------------------------------------------------
Every molecule renderer MUST pass ``props`` to all geometry helpers so that
individual card instances in ``deck.md`` can override values without touching
the theme.  The priority chain is: per-card prop → CSS token → computed default.

    # Canonical pattern — copy into every new molecule renderer
    pad         = ctx.card_pad_px(w, h, props)
    header_h    = ctx.card_header_h(w, h, props)
    header_gap  = ctx.card_header_gap(h, props)
    title_size  = ctx.card_header_font_size(title, text_w, h, props)
    icon_sz     = ctx.icon_size(w, h, props)
    icon_r      = ctx.icon_radius(icon_sz, props)
    hdr_color   = ctx.card_line_color("header", ctx.color("line-default"), props)
    ftr_color   = ctx.card_line_color("footer", ctx.color("line-default"), props)

Overridable props accepted in deck.md YAML front-matter (hyphen or underscore forms):

    # Geometry
    card-padding: 24
    card-header-height: 50
    card-header-gap: 12
    card-header-font-size: 14
    icon-size: 36
    icon-radius: 0

    # Colors
    header-line-color: "#E53935"
    footer-line-color: "#E53935"
    card_bg: filled | clean | alt | featured
    title-color: "#hex"
    body-color: "#hex"
    icon-bg: "#hex"
    icon-fg: "#hex"
    icon-stroke: "#hex"

    # Visibility  (none/false/0/off/hide → disabled; true/1/yes/on/show → enabled)
    show-header: none
    show-header-line: none
    show-footer: none
    show-footer-line: none

    # Alignment
    header-align: left | center | right
    header-line-width: 50%
    header-line-align: center
"""

from __future__ import annotations

import uuid
import html as _html
from xml.etree import ElementTree as ET


# ── Base ──────────────────────────────────────────────────────────────────────

class RenderContext:
    """Abstract base — defines the unified drawing API shared by all backends."""

    @staticmethod
    def is_transparent_fill(fill) -> bool:
        """Return True when *fill* represents no visible fill."""
        return str(fill or "").strip().lower() in {"", "none", "transparent", "rgba(0,0,0,0)"}

    # -- Primitives ------------------------------------------------------------

    def rect(self, x, y, w, h, fill, stroke=None, radius=0):
        raise NotImplementedError

    def text(self, x, y, w, h, text, size=14, bold=False, italic=False,
             color=None, align="left", valign="top", inner_margin: int = 0,
             font: str | None = None):
        """Draw a text box. inner_margin=0 removes default internal padding so
        text edges align precisely with adjacent shapes at the same x/y.
        font overrides the body font family for this call (e.g., for icon glyphs)."""
        raise NotImplementedError

    def badge(self, x, y, w, h, text, fill, text_color=None, size=11, radius=12):
        """Filled rounded pill with centered text."""
        self.rect(x, y, w, h, fill=fill, stroke=None, radius=radius)
        self.text(x, y, w, h, text=text, size=size, bold=True,
                  color=text_color or self.color("on-primary"), align="center", valign="middle")

    def divider(self, x, y, w, color=None):
        raise NotImplementedError

    def ellipse(self, x, y, w, h, fill, stroke=None):
        raise NotImplementedError

    def wedge(self, cx: int, cy: int, r: int,
              start_deg: float, end_deg: float,
              fill: str, hole_r: int = 0) -> None:
        """Filled pie sector from start_deg to end_deg (clockwise from 12 o'clock).

        hole_r > 0 produces a donut segment (annular sector).
        Subtended angle = end_deg - start_deg degrees.
        """
        raise NotImplementedError

    def line(self, x1, y1, x2, y2, color=None):
        raise NotImplementedError

    # -- Design tokens ---------------------------------------------------------

    @property
    def PAD(self) -> int:
        raise NotImplementedError

    @property
    def font(self) -> str:
        raise NotImplementedError

    def color(self, token: str) -> str:
        raise NotImplementedError

    def rad(self, key: str = "radius-medium") -> int:
        raise NotImplementedError

    def status_color(self, status: str):
        """Return (bg_hex, fg_hex) tuple for a semantic status string."""
        raise NotImplementedError

    def progress_color(self, pct: int) -> str:
        """Return a hex color indicating green/amber/red progress."""
        raise NotImplementedError

    # -- Card background helpers -----------------------------------------------

    _BG_VARIANT_MAP: dict[str, str] = {
        "default":  "bg-card",
        "clean":    "bg-card-clean",
        "filled":   "bg-card-filled",
        "alt":      "bg-card-alt",
        "featured": "bg-card-featured",
    }

    def card_bg_color(self, props: dict, default_token: str = "bg-card") -> str:
        """Return resolved hex for card background.

        Molecules call this instead of ``ctx.color("bg-card-xyz")`` so that
        deck.md can override any card's background via::

            card_bg: filled   # primary-container accent
            card_bg: clean    # plain surface white
            card_bg: featured # brand-primary vivid

        Valid override keys: default | clean | filled | alt | featured.
        Falls back to *default_token* when no override is present.
        """
        override = str(props.get("card_bg", "")).strip().lower()
        token = self._BG_VARIANT_MAP.get(override, default_token)
        return self.color(token)

    @staticmethod
    def _boolish(raw, default: bool = False) -> bool:
        """Interpret common user/theme truthy-falsy strings."""
        if raw is None:
            return default
        text = str(raw).strip().lower()
        if text in {"", "auto", "default"}:
            return default
        return text not in {"0", "false", "no", "off", "none", "hidden", "hide", "suppress"}

    def theme_var(self, name: str, default: str = "") -> str:
        """Read a resolved theme CSS custom property when available."""
        ds = getattr(self, "ds", None)
        sheet = getattr(ds, "stylesheet", None)
        if sheet is None:
            return default
        return sheet.var(name, default)

    def theme_flag(self, name: str, default: bool = False) -> bool:
        """Read a theme CSS custom property as a boolean-like flag."""
        return self._boolish(self.theme_var(name, "1" if default else "0"), default)

    @staticmethod
    def _prop_value(props: dict | None, *keys: str):
        if not isinstance(props, dict):
            return None
        for key in keys:
            if key in props:
                return props.get(key)
        return None

    def card_section_enabled(self, props: dict | None, section: str,
                             default: bool = True) -> bool:
        """Return whether a semantic card section should be rendered."""
        raw = self._prop_value(
            props,
            f"show_{section}",
            f"show-{section}",
            f"{section}_display",
            f"{section}-display",
        )
        if raw is not None:
            return self._boolish(raw, default)
        return self.theme_flag(f"--card-{section}-display", default)

    def card_line_enabled(self, props: dict | None, section: str,
                          default: bool = True) -> bool:
        """Return whether a semantic header/footer divider should be rendered."""
        raw = self._prop_value(
            props,
            f"show_{section}_line",
            f"show-{section}-line",
            f"{section}_line_display",
            f"{section}-line-display",
        )
        if raw is not None:
            return self._boolish(raw, default)
        return self.theme_flag(f"--card-{section}-line-display", default)

    def card_line_color(self, section: str, default_color: str,
                        props: dict | None = None) -> str:
        """Return the shared divider color for a semantic card section.

        Per-card override: set ``header-line-color: #hex`` or
        ``footer-line-color: #hex`` in the card's deck.md props block.
        Theme-wide: ``--card-{section}-line-color`` CSS token.
        """
        raw = self._prop_value(props, f"{section}-line-color", f"{section}_line_color")
        if raw:
            return str(raw)
        return self.theme_var(f"--card-{section}-line-color", default_color)

    @staticmethod
    def _ratioish(raw, default: float = 1.0) -> float:
        if raw is None:
            return default
        text = str(raw).strip().lower()
        if not text:
            return default
        try:
            if text.endswith("%"):
                return max(0.0, min(1.0, float(text[:-1].strip()) / 100.0))
            value = float(text)
            if value > 1.0:
                return max(0.0, min(1.0, value / 100.0))
            return max(0.0, min(1.0, value))
        except Exception:
            return default

    def card_title_color(self, props: dict | None = None,
                         default_token: str = "text-on-muted") -> str:
        raw = self._prop_value(props, "card_title_color", "card-title-color", "title_color", "title-color")
        if raw:
            return str(raw)
        return self.theme_var("--card-title-color", self.color(default_token))

    def card_body_color(self, props: dict | None = None,
                        default_token: str = "text-secondary") -> str:
        raw = self._prop_value(props, "card_body_color", "card-body-color", "body_color", "body-color")
        if raw:
            return str(raw)
        return self.theme_var("--card-body-color", self.color(default_token))

    def card_subtitle_color(self, props: dict | None = None,
                            default_token: str = "text-secondary") -> str:
        raw = self._prop_value(props, "card_subtitle_color", "card-subtitle-color", "subtitle_color", "subtitle-color")
        if raw:
            return str(raw)
        return self.theme_var("--card-subtitle-color", self.color(default_token))

    def card_header_align(self, props: dict | None = None, default: str = "left") -> str:
        raw = self._prop_value(props, "header_align", "header-align", "title_align", "title-align")
        align = str(raw if raw is not None else self.theme_var("--card-header-title-align", default)).strip().lower()
        return align if align in {"left", "center", "right"} else default

    def card_icon_align(self, props: dict | None = None, default: str = "right") -> str:
        raw = self._prop_value(props, "icon_align", "icon-align", "header_icon_align", "header-icon-align")
        align = str(raw if raw is not None else self.theme_var("--card-header-icon-align", default)).strip().lower()
        return align if align in {"left", "right"} else default

    # ── Card geometry helpers ─────────────────────────────────────────────
    # Centralised sizing so all molecule renderers produce consistent
    # header baselines on cards that share the same row (same height).
    # These read optional CSS tokens to allow theme-wide overrides.
    #
    # Templates set ``ctx.ref_h = content_h`` before their dispatch loop so
    # that all cards on the same slide (regardless of individual height)
    # receive identical header heights, icon sizes, and gaps.
    # Templates clear it with ``ctx.ref_h = None`` after the loop.

    def _ref(self, h: int) -> int:
        """Return the slide reference height when set, else the card's own h."""
        ref = getattr(self, "ref_h", None)
        return ref if (ref and ref > 0) else h

    def card_pad_px(self, w: int, h: int, props: dict | None = None) -> int:
        """Inner card padding in px.

        Uses **card height only** (not ``min(w, h)``) so all cards in the
        same grid row share the same padding regardless of column width.
        Per-card override: ``card-padding: <px>`` in the card's deck.md props block.
        Theme-wide: ``--card-padding`` CSS token.
        """
        raw_prop = self._prop_value(props, "card-padding", "card_padding")
        if raw_prop is not None:
            try:
                v = int(float(str(raw_prop)))
                if v > 0:
                    return v
            except (ValueError, TypeError):
                pass
        raw = self.theme_var("--card-padding", "").strip()
        try:
            v = int(float(raw)) if raw else 0
            if v > 0:
                return v
        except (ValueError, TypeError):
            pass
        return max(self.PAD, int(h * 0.055))

    def card_header_h(self, w: int, h: int, props: dict | None = None) -> int:
        """Standard card header height in px.

        Shared by all molecule renderers so the divider line sits at the
        same Y position across every card in a grid row.  The icon badge
        fits inside this height; it does **not** inflate it.
        Uses ``ctx.ref_h`` (slide content height) when available so that
        cards of different heights on the same slide get identical header zones.
        Per-card override: ``card-header-height: <px>`` in the card's deck.md props block.
        Theme-wide: ``--card-header-height`` CSS token.
        """
        raw_prop = self._prop_value(props, "card-header-height", "card_header_height")
        if raw_prop is not None:
            try:
                v = int(float(str(raw_prop)))
                if v > 0:
                    return v
            except (ValueError, TypeError):
                pass
        raw = self.theme_var("--card-header-height", "").strip()
        try:
            v = int(float(raw)) if raw else 0
            if v > 0:
                return v
        except (ValueError, TypeError):
            pass
        return max(34, int(self._ref(h) * 0.12))

    def card_header_gap(self, h: int, props: dict | None = None) -> int:
        """Gap between the header row and the divider line in px.

        Uses ``ctx.ref_h`` (slide content height) when available so that
        cards of different heights on the same slide get the identical gap.
        Per-card override: ``card-header-gap: <px>`` in the card's deck.md props block.
        Theme-wide: ``--card-header-gap`` CSS token.
        """
        raw_prop = self._prop_value(props, "card-header-gap", "card_header_gap")
        if raw_prop is not None:
            try:
                v = int(float(str(raw_prop)))
                if v > 0:
                    return v
            except (ValueError, TypeError):
                pass
        raw = self.theme_var("--card-header-gap", "").strip()
        try:
            v = int(float(raw)) if raw else 0
            if v > 0:
                return v
        except (ValueError, TypeError):
            pass
        return max(self.spacing("s"), int(self._ref(h) * 0.018))

    def card_header_font_size(self, title: str = "", text_w: int = 200,
                               h: int = 300,
                               props: dict | None = None) -> int:
        """Header / title font size in pt.

        Uses ``ctx.ref_h`` (slide content height) when available.
        Per-card override: ``card-header-font-size: <pt>`` in the card's deck.md props block.
        Theme-wide: ``--card-header-font-size`` CSS token.
        """
        raw_prop = self._prop_value(props, "card-header-font-size", "card_header_font_size")
        if raw_prop is not None:
            try:
                v = int(float(str(raw_prop)))
                if v > 0:
                    return v
            except (ValueError, TypeError):
                pass
        raw = self.theme_var("--card-header-font-size", "").strip()
        try:
            v = int(float(raw)) if raw else 0
            if v > 0:
                return v
        except (ValueError, TypeError):
            pass
        ref = self._ref(h)
        cap = max(self.font_size("body"), min(22, int(ref * 0.034)))
        if not title or text_w <= 0:
            return cap
        return self.fit_text_size(
            title, text_w,
            max_size=cap,
            min_size=self.font_size("label"),
            bold=True, safety=0.92,
        )

    # ── Icon badge helpers ────────────────────────────────────────────────
    # These read the canonical --color-icon-* / --icon-* CSS tokens so every
    # molecule gets a consistent single source of truth.  Pass a props dict
    # to allow per-card overrides via deck.md front-matter.

    def icon_bg(self, props: dict | None = None,
                default_token: str = "primary-container") -> str:
        """Badge fill color — reads ``--color-icon-bg`` CSS token."""
        raw = self._prop_value(props, "icon_bg", "icon-bg")
        if raw:
            return str(raw)
        return self.theme_var("--color-icon-bg", self.color(default_token))

    def icon_fg(self, props: dict | None = None,
                default_token: str = "primary") -> str:
        """Badge glyph/icon color — reads ``--color-icon-fg`` CSS token."""
        raw = self._prop_value(props, "icon_fg", "icon-fg")
        if raw:
            return str(raw)
        return self.theme_var("--color-icon-fg", self.color(default_token))

    def icon_stroke(self, props: dict | None = None,
                    default_token: str = "border-default") -> str:
        """Badge border/outline color — reads ``--color-icon-stroke`` CSS token."""
        raw = self._prop_value(props, "icon_stroke", "icon-stroke")
        if raw:
            return str(raw)
        return self.theme_var("--color-icon-stroke", self.color(default_token))

    def icon_size(self, card_w: int, card_h: int,
                  props: dict | None = None) -> int:
        """Compute icon-badge side length in px.

        Uses ``--icon-size`` CSS token as the base value when set; otherwise
        uses ``ctx.ref_h`` (slide content height) so all cards on the same
        slide get a consistent badge size regardless of individual card height.
        Result is always clamped to [24, 96] px.
        Per-card override: ``icon-size: <px>`` in the card's deck.md props block.
        """
        raw_prop = self._prop_value(props, "icon-size", "icon_size")
        if raw_prop is not None:
            try:
                v = int(float(str(raw_prop)))
                if v > 0:
                    return max(24, min(96, v))
            except (ValueError, TypeError):
                pass
        raw = self.theme_var("--icon-size", "").strip()
        try:
            base = int(float(raw)) if raw else None
        except (ValueError, TypeError):
            base = None
        if base and base > 0:
            return max(24, min(96, base))
        ref = self._ref(card_h)
        return max(36, min(72, int(ref * 0.14)))

    def icon_radius(self, size: int, props: dict | None = None) -> int:
        """Compute icon-badge corner radius in px.

        Uses ``--icon-radius`` CSS token when set; otherwise derives the radius
        from the badge size and the ``radius-sharp`` / ``radius-large`` design
        tokens.
        Per-card override: ``icon-radius: <px>`` in the card's deck.md props block.
        """
        raw_prop = self._prop_value(props, "icon-radius", "icon_radius")
        if raw_prop is not None:
            try:
                r = int(float(str(raw_prop)))
                return max(0, r)
            except (ValueError, TypeError):
                pass
        raw = self.theme_var("--icon-radius", "").strip()
        try:
            r = int(float(raw)) if raw else None
        except (ValueError, TypeError):
            r = None
        if r is not None:
            return max(0, r)
        return max(self.rad("radius-sharp"), min(size // 3, self.rad("radius-large")))

    def card_divider_span(self, section: str, x: int, w: int,
                          props: dict | None = None) -> tuple[int, int]:
        align_raw = self._prop_value(props, f"{section}_line_align", f"{section}-line-align")
        width_raw = self._prop_value(props, f"{section}_line_width", f"{section}-line-width")
        align = str(align_raw if align_raw is not None else self.theme_var(f"--card-{section}-line-align", "left")).strip().lower()
        if align not in {"left", "center", "right"}:
            align = "left"
        ratio = self._ratioish(width_raw if width_raw is not None else self.theme_var(f"--card-{section}-line-width", "100%"), 1.0)
        span_w = max(1, int(w * ratio))
        if align == "center":
            start_x = x + max(0, (w - span_w) // 2)
        elif align == "right":
            start_x = x + max(0, w - span_w)
        else:
            start_x = x
        return start_x, span_w

    @staticmethod
    def estimate_text_width(text: str, size: int | float, bold: bool = False) -> float:
        """Approximate single-line text width in px for responsive fitting."""
        lines = str(text or "").splitlines() or [""]
        widest = 0.0
        for line in lines:
            units = 0.0
            for ch in line:
                if ch.isspace():
                    units += 0.33
                elif ch in "ilI.,'`:;!|[](){}":
                    units += 0.32
                elif ch in "MW@#%&€$":
                    units += 0.92
                elif ch.isdigit():
                    units += 0.62
                elif ch.isupper():
                    units += 0.70
                else:
                    units += 0.56
            widest = max(widest, units)
        return widest * float(size) * (1.04 if bold else 1.0)

    def fit_text_size(self, text: str, max_width: int | float, max_size: int,
                      min_size: int | None = None, bold: bool = False,
                      safety: float = 0.92) -> int:
        """Shrink text size until its estimated width fits within *max_width*."""
        floor = max(1, int(min_size if min_size is not None else self.font_size("caption")))
        size = max(floor, int(max_size))
        allowed = max(1.0, float(max_width) * safety)
        while size > floor and self.estimate_text_width(text, size, bold=bold) > allowed:
            size -= 1
        return max(floor, size)

    @property
    def gutter(self) -> int:
        """Card/column gap in px (canvas.gutter from design-config)."""
        raise NotImplementedError

    def spacing(self, name: str = "m") -> int:
        """Return spacing token value in px (xs/s/m/l/xl)."""
        raise NotImplementedError

    def font_size(self, role: str = "body") -> int:
        """Return typography size for *role* (body/heading/heading-sub/caption…)."""
        raise NotImplementedError

    def font_bold(self, role: str = "body") -> bool:
        """Return True when the typography *role* weight is >= 600 (Semibold+)."""
        raise NotImplementedError

    # -- Icon-set helpers -------------------------------------------------------

    def icon_set(self) -> str:
        """Return the active icon set name from ``--icon-set`` CSS token.

        Built-in values: ``emoji`` (default) | ``noto`` | ``material`` |
        ``material-outlined`` | ``material-rounded`` | ``material-sharp`` |
        ``material-two-tone``.
        """
        return self.theme_var("--icon-set", "emoji").strip().lower() or "emoji"

    def icon_font(self) -> str | None:
        """Return the icon font family override from ``--icon-font-family``.

        When the token is empty (default), returns ``None`` so ``ctx.text()``
        falls back to the body font.  When set to e.g. ``"Material Icons Outlined"``
        the text call uses that font, enabling Material Icons ligature rendering.
        """
        raw = self.theme_var("--icon-font-family", "").strip()
        return raw if raw else None

    def resolve_icon(self, description: str, prefer_letter: bool = False) -> str:
        """Resolve a natural-language icon description using the active icon set.

        Delegates to ``input_utils.resolve_icon`` with the current ``icon_set()``.
        Use this in card renderers instead of importing ``resolve_icon`` directly
        so the CSS ``--icon-set`` token is always respected.
        """
        from rendering.input_utils import resolve_icon as _ri
        iset = self.icon_set()
        # For ligature-based sets never prefer single letters
        force_letter = prefer_letter and iset in ("emoji", "noto", "")
        return _ri(description, prefer_letter=force_letter, icon_set=iset)

    def draw_icon(self, x: int, y: int, w: int, h: int,
                  icon_name: str, color: str | None = None) -> None:
        """Draw an icon centred in the bounding box (x, y, w, h).

        Subclasses override this to embed a PNG rendered from SVG path data.
        The base implementation falls back to a single-letter or emoji glyph
        via ``self.text()`` (no special icon font required).
        """
        # Resolve concept → emoji/letter (prefer_letter=True avoids font dep)
        glyph = self.resolve_icon(icon_name, prefer_letter=True)
        self.text(x, y, w, h, glyph,
                  size=int(min(w, h) * 0.65) or 14,
                  bold=False,
                  color=color or self.color("on-surface"),
                  align="center", valign="middle")

    # ── Layout chrome helpers ─────────────────────────────────────────────────
    # Read --content-block-gap / --content-area-padding / --content-block-padding
    # from the CSS theme so layout templates are fully token-driven.

    @property
    def content_gap(self) -> int:
        """Gap between adjacent content blocks in px — reads ``--content-block-gap``.
        Falls back to ``gutter`` when the token is absent or zero."""
        raw = self.theme_var("--content-block-gap", "").strip()
        try:
            v = int(float(raw)) if raw else 0
            return v if v > 0 else self.gutter
        except (ValueError, TypeError):
            return self.gutter

    @property
    def content_padding(self) -> int:
        """Inner padding of the content area in px — reads ``--content-area-padding``."""
        raw = self.theme_var("--content-area-padding", "0").strip()
        try:
            return max(0, int(float(raw)))
        except (ValueError, TypeError):
            return 0

    @property
    def block_padding(self) -> int:
        """Inner padding of each content block in px — reads ``--content-block-padding``."""
        raw = self.theme_var("--content-block-padding", "0").strip()
        try:
            return max(0, int(float(raw)))
        except (ValueError, TypeError):
            return 0

    @property
    def block_bg_color(self) -> str:
        """Default block fill color — reads ``--content-block-bg-color``."""
        return self.theme_var("--content-block-bg-color", "transparent")


# ── draw.io backend ───────────────────────────────────────────────────────────

class DrawioCtx(RenderContext):
    """Renders primitives by appending mxCell elements to an XML root node."""

    def __init__(self, root, ds):
        """
        Parameters
        ----------
        root : xml.etree.ElementTree.Element
            The <root> element of the diagram's mxGraphModel.
        ds : DesignSystem
            Loaded design-system instance.
        """
        self.root = root
        self.ds   = ds

    # -- Helpers ---------------------------------------------------------------

    @staticmethod
    def _new_id() -> str:
        return uuid.uuid4().hex[:16]

    # -- Primitives ------------------------------------------------------------

    def rect(self, x, y, w, h, fill, stroke=None, radius=0, value="", extra_style=""):
        rounded  = "1" if radius > 0 else "0"
        arc_size = max(1, int(radius))
        fill_color = "none" if self.is_transparent_fill(fill) else fill
        fill_opacity = "0" if self.is_transparent_fill(fill) else "100"
        style = (f"rounded={rounded};arcSize={arc_size};"
                 f"fillColor={fill_color};fillOpacity={fill_opacity};strokeColor={stroke or 'none'};"
                 f"whiteSpace=wrap;html=1;{extra_style}")
        cell = ET.SubElement(self.root, "mxCell",
                             {"id": self._new_id(), "parent": "1",
                              "value": _html.escape(str(value)) if value else "",
                              "style": style, "vertex": "1"})
        ET.SubElement(cell, "mxGeometry",
                      {"x": str(int(x)), "y": str(int(y)),
                       "width": str(int(w)), "height": str(int(h)),
                       "as": "geometry"})
        return cell

    def text(self, x, y, w, h, text, size=14, bold=False, italic=False,
             color=None, align="left", valign="top", inner_margin: int = 0,
             font: str | None = None):
        _font = font if font else self.font
        font_style = (1 if bold else 0) | (2 if italic else 0)
        margin_css  = (f"spacing={inner_margin};spacingLeft={inner_margin};"
                       f"spacingRight={inner_margin};spacingTop={inner_margin};"
                       f"spacingBottom={inner_margin};" if inner_margin else "")
        style = (f"text;html=1;align={align};verticalAlign={valign};"
                 f"fontSize={size};fontStyle={font_style};"
                 f"fontFamily={_font};"
                 f"fontColor={color or self.color('on-surface')};"
                 f"whiteSpace=wrap;overflow=hidden;{margin_css}")
        cell = ET.SubElement(self.root, "mxCell",
                             {"id": self._new_id(), "parent": "1",
                              "value": _html.escape(str(text)),
                              "style": style, "vertex": "1"})
        ET.SubElement(cell, "mxGeometry",
                      {"x": str(int(x)), "y": str(int(y)),
                       "width": str(int(w)), "height": str(int(h)),
                       "as": "geometry"})
        return cell

    def draw_icon(self, x: int, y: int, w: int, h: int,
                  icon_name: str, color: str | None = None) -> None:
        """Embed a Material Icon as a base64-encoded PNG image cell.

        Falls back to the text glyph if SVG path data is unavailable.
        """
        import base64
        from rendering.input_utils import icon_to_png_bytes, resolve_icon as _ri
        sz   = int(min(w, h))
        col  = color or self.color("on-surface")
        iset = self.icon_set()
        # Resolve concept name → Material ligature (e.g. "employees" → "group")
        lig  = _ri(icon_name, prefer_letter=False, icon_set=iset)
        png  = icon_to_png_bytes(lig, sz, col)
        if png:
            b64   = base64.b64encode(png).decode()
            style = (f"shape=image;imageAspect=1;aspect=fixed;"
                     f"image=data:image/png;base64,{b64};")
            cell = ET.SubElement(self.root, "mxCell",
                                 {"id": self._new_id(), "parent": "1",
                                  "value": "", "style": style, "vertex": "1"})
            ET.SubElement(cell, "mxGeometry",
                          {"x": str(int(x)), "y": str(int(y)),
                           "width": str(int(w)), "height": str(int(h)),
                           "as": "geometry"})
        else:
            # Fallback: single-letter badge — never render a multi-char ligature
            # name as text (it would wrap letter-by-letter in a small box).
            glyph = _ri(icon_name, prefer_letter=True, icon_set=iset)
            if len(glyph) > 2:  # still a ligature string, not a char/emoji
                glyph = (icon_name.strip()[0].upper()
                         if icon_name.strip() else "?")
            self.text(x, y, w, h, glyph,
                      size=int(sz * 0.65) or 14,
                      color=col, align="center", valign="middle")

    def divider(self, x, y, w, color=None):
        c     = color or self.color("border-subtle")
        style = f"endArrow=none;html=1;strokeColor={c};strokeWidth=1;"
        cell  = ET.SubElement(self.root, "mxCell",
                              {"id": self._new_id(), "parent": "1",
                               "value": "", "style": style, "edge": "1"})
        geo   = ET.SubElement(cell, "mxGeometry",
                              {"relative": "1", "as": "geometry"})
        ET.SubElement(geo, "mxPoint",
                      {"x": str(int(x)), "y": str(int(y)), "as": "sourcePoint"})
        ET.SubElement(geo, "mxPoint",
                      {"x": str(int(x + w)), "y": str(int(y)), "as": "targetPoint"})

    def line(self, x1, y1, x2, y2, color=None):
        """Alias for divider supporting arbitrary start/end points."""
        c     = color or self.color("border-subtle")
        style = f"endArrow=none;html=1;strokeColor={c};strokeWidth=1;"
        cell  = ET.SubElement(self.root, "mxCell",
                              {"id": self._new_id(), "parent": "1",
                               "value": "", "style": style, "edge": "1"})
        geo   = ET.SubElement(cell, "mxGeometry",
                              {"relative": "1", "as": "geometry"})
        ET.SubElement(geo, "mxPoint",
                      {"x": str(int(x1)), "y": str(int(y1)), "as": "sourcePoint"})
        ET.SubElement(geo, "mxPoint",
                      {"x": str(int(x2)), "y": str(int(y2)), "as": "targetPoint"})

    def ellipse(self, x, y, w, h, fill, stroke=None):
        fill_color = "none" if self.is_transparent_fill(fill) else fill
        fill_opacity = "0" if self.is_transparent_fill(fill) else "100"
        style = (f"ellipse;fillColor={fill_color};fillOpacity={fill_opacity};"
                 f"strokeColor={stroke or 'none'};html=1;")
        cell  = ET.SubElement(self.root, "mxCell",
                              {"id": self._new_id(), "parent": "1",
                               "value": "", "style": style, "vertex": "1"})
        ET.SubElement(cell, "mxGeometry",
                      {"x": str(int(x)), "y": str(int(y)),
                       "width": str(int(w)), "height": str(int(h)),
                       "as": "geometry"})

    def wedge(self, cx: int, cy: int, r: int,
              start_deg: float, end_deg: float,
              fill: str, hole_r: int = 0) -> None:
        """Filled pie sector using DrawIO arc shape (0° = 12 o\'clock, clockwise).

        DrawIO style=arc uses 0° = 3 o\'clock, clockwise.  We convert our convention
        (start from top) by subtracting 90°.
        For donut segments the caller overlays a center circle after all sectors.
        """
        dw_start = (start_deg - 90) % 360
        dw_end   = (end_deg   - 90) % 360
        style = (f"arc;startAngle={dw_start:.2f};endAngle={dw_end:.2f};"
                 f"fillColor={fill};strokeColor=none;html=1;")
        cell = ET.SubElement(self.root, "mxCell",
                             {"id": self._new_id(), "parent": "1",
                              "value": "", "style": style, "vertex": "1"})
        ET.SubElement(cell, "mxGeometry",
                      {"x": str(int(cx - r)), "y": str(int(cy - r)),
                       "width": str(int(r * 2)), "height": str(int(r * 2)),
                       "as": "geometry"})

    # -- Design tokens (resolved from CSS) -------------------------------------

    @property
    def PAD(self) -> int:
        return int(self.ds.canvas().get("padding", 20))

    @property
    def font(self) -> str:
        return self.ds.font_family()

    def color(self, token: str) -> str:
        return self.ds.color(token)

    def rad(self, key: str = "radius-medium") -> int:
        return self.ds.radius(key)

    def status_color(self, status: str):
        m = {
            "success": (self.color("success"),         self.color("on-success")),
            "warning": (self.color("warning"),          self.color("on-surface")),
            "error":   (self.color("error"),            self.color("on-error")),
            "neutral": (self.color("surface-variant"),  self.color("on-surface-variant")),
        }
        return m.get(status, m["neutral"])

    def progress_color(self, pct: int) -> str:
        if pct >= 80: return self.color("success")
        if pct >= 40: return self.color("warning")
        return self.color("error")

    @property
    def gutter(self) -> int:
        return int(self.ds.canvas().get("gutter", 24))

    def spacing(self, name: str = "m") -> int:
        return self.ds.spacing(name)

    def font_size(self, role: str = "body") -> int:
        return self.ds.font_size(role)

    def font_bold(self, role: str = "body") -> bool:
        return self.ds.font_bold(role)


# ── PowerPoint backend ────────────────────────────────────────────────────────

class PptxCtx(RenderContext):
    """Renders primitives by adding python-pptx shapes to a slide object."""

    def __init__(self, slide, ds):
        """
        Parameters
        ----------
        slide : pptx.slide.Slide
            python-pptx slide object to add shapes to.
        ds : DesignSystem
            Loaded design-system instance.
        """
        self.slide = slide
        self.ds    = ds

    # -- Helpers ---------------------------------------------------------------

    @staticmethod
    def _px(v) -> int:
        """Pixels → EMU (96 dpi: 1 px = 9525 EMU)."""
        return int(v * 9525)

    @staticmethod
    def _rgb(hex_color):
        from pptx.dml.color import RGBColor
        if not hex_color:
            return None
        h = str(hex_color).strip().lower()
        if h in {"", "none", "transparent", "rgba(0,0,0,0)"}:
            return None
        if h.startswith("#"):
            h = h.lstrip("#")
        if len(h) == 3:
            h = "".join(c * 2 for c in h)
        if len(h) != 6:
            raise ValueError(f"Invalid color value: {hex_color!r}")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    @staticmethod
    def _apply_radius(shape, radius_pt: float):
        """Apply rounded-rect corner radius via XML manipulation."""
        from pptx.oxml.ns import qn
        from lxml import etree  # type: ignore[attr-defined]
        sp     = shape.element
        spPr   = sp.find(qn("p:spPr"))
        if spPr is None:
            return
        prstGeom = spPr.find(qn("a:prstGeom"))
        if prstGeom is None:
            return
        prstGeom.set("prst", "roundRect")
        avLst = prstGeom.find(qn("a:avLst"))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn("a:avLst"))
        avLst.clear()
        gd = etree.SubElement(avLst, qn("a:gd"))
        gd.set("name", "adj")
        gd.set("fmla", f"val {min(50000, max(0, int(radius_pt * 1000)))}")

    # -- Primitives ------------------------------------------------------------

    def rect(self, x, y, w, h, fill, stroke=None, radius=0, **_):
        from pptx.util import Emu
        px = self._px
        shape = self.slide.shapes.add_shape(
            1, Emu(px(x)), Emu(px(y)), Emu(px(w)), Emu(px(h)))
        if self.is_transparent_fill(fill):
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._rgb(fill)
        if stroke and not self.is_transparent_fill(stroke):
            color = self._rgb(stroke)
            if color is not None:
                shape.line.color.rgb = color
                shape.line.width = Emu(px(1))
            else:
                shape.line.fill.background()
        else:
            shape.line.fill.background()
        if radius > 0:
            self._apply_radius(shape, radius)
        if self._shadow_level() == 0:
            self._disable_shadow(shape)
        return shape

    def text(self, x, y, w, h, text, size=14, bold=False, italic=False,
             color=None, align="left", valign="top", inner_margin: int = 0,
             font: str | None = None):
        from pptx.util import Emu, Pt
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        _font = font if font else self.font
        px = self._px
        box = self.slide.shapes.add_textbox(Emu(px(x)), Emu(px(y)), Emu(px(w)), Emu(px(h)))
        tf  = box.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        _m = Emu(px(inner_margin))
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = _m
        va_map = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                  "bottom": MSO_ANCHOR.BOTTOM}
        try:
            tf.vertical_anchor = va_map.get(valign, MSO_ANCHOR.TOP)
        except Exception:
            pass
        p  = tf.paragraphs[0]
        al_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                  "right": PP_ALIGN.RIGHT}
        p.alignment = al_map.get(align, PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = str(text)
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.italic = italic
        run.font.color.rgb = self._rgb(color or self.color("on-surface"))
        run.font.name  = _font
        return box

    def draw_icon(self, x: int, y: int, w: int, h: int,
                  icon_name: str, color: str | None = None) -> None:
        """Embed a Material Icon as a PNG picture in the slide.

        The PNG is rendered via cairosvg from the built-in SVG path data,
        sized to ``min(w, h)`` pixels and centred inside the bounding box.
        Falls back to an emoji/letter glyph when path data is unavailable.
        """
        from io import BytesIO
        from pptx.util import Emu
        from rendering.input_utils import icon_to_png_bytes, resolve_icon as _ri
        px   = self._px
        sz   = int(min(w, h))
        col  = color or self.color("on-surface")
        iset = self.icon_set()
        # Resolve concept → Material ligature (e.g. "employees" → "group")
        lig  = _ri(icon_name, prefer_letter=False, icon_set=iset)
        png  = icon_to_png_bytes(lig, sz, col)
        if png:
            # Centre the square PNG inside the (possibly non-square) box
            cx = x + (w - sz) // 2
            cy = y + (h - sz) // 2
            self.slide.shapes.add_picture(
                BytesIO(png),
                Emu(px(cx)), Emu(px(cy)),
                Emu(px(sz)), Emu(px(sz)),
            )
        else:
            # Fallback: single-letter badge — never render a multi-char ligature
            # name as text (it would wrap letter-by-letter in a small box).
            glyph = _ri(icon_name, prefer_letter=True, icon_set=iset)
            if len(glyph) > 2:  # still a ligature string, not a char/emoji
                glyph = (icon_name.strip()[0].upper()
                         if icon_name.strip() else "?")
            self.text(x, y, w, h, glyph,
                      size=int(sz * 0.65) or 14,
                      color=col, align="center", valign="middle")

    def divider(self, x, y, w, color=None):
        self.line(x, y, x + w, y, color)

    def line(self, x1, y1, x2, y2, color=None):
        from pptx.util import Emu
        px  = self._px
        c   = color or self.color("border-subtle")
        if y1 == y2:  # horizontal
            shape = self.slide.shapes.add_shape(
                1, Emu(px(x1)), Emu(px(y1)), Emu(px(x2 - x1)), Emu(px(1)))
        else:         # vertical
            shape = self.slide.shapes.add_shape(
                1, Emu(px(x1)), Emu(px(y1)), Emu(px(1)), Emu(px(y2 - y1)))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(c)
        shape.line.fill.background()
        if self._shadow_level() == 0:
            self._disable_shadow(shape)

    @staticmethod
    def _disable_shadow(shape) -> None:
        """Suppress the PowerPoint theme shadow by zeroing the effectRef index.

        Shapes rendered via ``add_shape()`` inherit ``<a:effectRef idx="2">``
        from the slide theme, which produces a drop shadow.  Setting idx to 0
        removes the theme-effect reference entirely.
        """
        from pptx.oxml.ns import qn
        sp    = shape.element
        style = sp.find(qn("p:style"))
        if style is None:
            return
        effect_ref = style.find(qn("a:effectRef"))
        if effect_ref is not None:
            effect_ref.set("idx", "0")

    def _shadow_level(self) -> int:
        """Read --shape-shadow-level from the design system (0 = no shadow)."""
        raw = self.ds.stylesheet.var("--shape-shadow-level", "0")
        try:
            return int(float(raw))
        except (ValueError, TypeError):
            return 0

    def ellipse(self, x, y, w, h, fill, stroke=None):
        from pptx.util import Emu
        px    = self._px
        shape = self.slide.shapes.add_shape(
            9, Emu(px(x)), Emu(px(y)), Emu(px(w)), Emu(px(h)))
        if self.is_transparent_fill(fill):
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._rgb(fill)
        if stroke and not self.is_transparent_fill(stroke):
            color = self._rgb(stroke)
            if color is not None:
                shape.line.color.rgb = color
                shape.line.width = Emu(px(2))
            else:
                shape.line.fill.background()
        else:
            shape.line.fill.background()

    def wedge(self, cx: int, cy: int, r: int,
              start_deg: float, end_deg: float,
              fill: str, hole_r: int = 0) -> None:
        """Filled pie sector via freeform polygon (0° = 12 o\'clock, clockwise).

        Builds a polygon approximation with ~3° resolution arc steps.
        hole_r > 0 draws an annular (donut) segment.
        """
        import math

        # Convert from "clockwise from top" to standard math angles (radians)
        def _to_rad(cw_deg: float) -> float:
            return math.radians(cw_deg - 90)

        n_steps = max(6, int(abs(end_deg - start_deg) // 3))
        angles  = [start_deg + (end_deg - start_deg) * i / n_steps
                   for i in range(n_steps + 1)]

        outer_pts = [(cx + r * math.cos(_to_rad(a)),
                      cy + r * math.sin(_to_rad(a))) for a in angles]

        if hole_r > 0:
            inner_pts = [(cx + hole_r * math.cos(_to_rad(a)),
                          cy + hole_r * math.sin(_to_rad(a)))
                         for a in reversed(angles)]
            poly = outer_pts + inner_pts
        else:
            poly = [(cx, cy)] + outer_pts

        start_x, start_y = poly[0]
        remaining = poly[1:]

        try:
            # pptx >= 1.0 API: build_freeform(x, y, scale) + add_line_segments()
            # scale=9525 → 1 canvas-px = 9525 EMU
            fb = self.slide.shapes.build_freeform(
                int(start_x), int(start_y), scale=9525)
            fb.add_line_segments(
                [(int(x2), int(y2)) for x2, y2 in remaining], close=True)
            shape = fb.convert_to_shape()
        except AttributeError:
            # Fallback: simple ellipse for older / incompatible pptx builds
            self.ellipse(cx - r, cy - r, r * 2, r * 2, fill=fill)
            return

        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(fill)
        shape.line.fill.background()

    # -- Design tokens (resolved from CSS) ------------------------------------

    @property
    def PAD(self) -> int:
        return int(self.ds.canvas().get("padding", 20))

    @property
    def font(self) -> str:
        return self.ds.font_family()

    def color(self, token: str) -> str:
        return self.ds.color(token)

    def rad(self, key: str = "radius-medium") -> int:
        return self.ds.radius(key)

    def status_color(self, status: str):
        m = {
            "success": (self.color("success"),         self.color("on-success")),
            "warning": (self.color("warning"),          self.color("on-surface")),
            "error":   (self.color("error"),            self.color("on-error")),
            "neutral": (self.color("surface-variant"),  self.color("on-surface-variant")),
        }
        return m.get(status, m["neutral"])

    def progress_color(self, pct: int) -> str:
        if pct >= 80: return self.color("success")
        if pct >= 40: return self.color("warning")
        return self.color("error")

    @property
    def gutter(self) -> int:
        return int(self.ds.canvas().get("gutter", 24))

    def spacing(self, name: str = "m") -> int:
        return self.ds.spacing(name)

    def font_size(self, role: str = "body") -> int:
        return self.ds.font_size(role)

    def font_bold(self, role: str = "body") -> bool:
        return self.ds.font_bold(role)
