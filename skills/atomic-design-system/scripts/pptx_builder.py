"""
pptx_builder.py — Generate PowerPoint (.pptx) from a resolved slide plan
─────────────────────────────────────────────────────────────────────────────
Architecture:

    PptxBuilder            ← thin orchestrator: scaffolds slides, delegates layout
        build()            ← public entry-point
        _build_slide()     ← scaffold (background, title, divider) + TEMPLATE_REGISTRY lookup
        _background()      ← full-slide background rect
        _slide_title()     ← heading + divider rule, returns content_y
        _make_slide()      ← creates a blank python-pptx slide
        _canvas()          ← returns (width, height, margin) from design system

    TEMPLATE_REGISTRY      ← maps layout slug → layout renderer (rendering/templates/)
        HeroTitleLayout    ← hero-title
        GridLayout         ← grid-3 / comparison-2col
        DataInsightLayout  ← data-insight
        NumberedListLayout ← numbered-list

    slide_helpers          ← shared parsing + dispatch (used by both builders)
        _extract_section_blocks()
        dispatch()

Usage:
    python scripts/pptx_builder.py deck.md [--stylesheet theme.css] [--output deck.pptx]

Requires: python-pptx, pyyaml
"""

from __future__ import annotations
import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from lxml import etree

sys.path.insert(0, str(Path(__file__).parent))
from design_system_utils import DesignSystem
from md_parser import parse_markdown, Slide
from agenda_injector import inject_agenda_slides, materialize_agenda_to_deck, strip_agenda_from_deck
from rendering.context import PptxCtx
from rendering.templates import TEMPLATE_REGISTRY
from slide_helpers import _extract_section_blocks, dispatch
from qa_render import render as qa_render


# ── PX helper (pixels → EMU at 96 dpi: 1 px = 914400 ÷ 96 = 9525 EMU) ───────

def px(v: int | float) -> int:
    return int(v * 9525)


# ── PowerPoint section injection ──────────────────────────────────────────────

def _inject_pptx_sections(prs: Presentation, slides: list) -> None:
    """Inject p14:sectionLst into the presentation XML so PowerPoint shows
    section separators in the slide panel and presenter view.

    Sections are derived from ``slide.section`` on each Slide object.
    Slides with ``section == None`` or ``section_index == -1`` are placed
    in an implicit first section named after the first slide's title.

    Only runs when at least 2 distinct section names are present.
    """
    import uuid as _uuid

    # Collect ordered (section_name → [pptx_slide_index]) mapping
    section_order: list[str] = []
    section_map: dict[str, list[int]] = {}

    # Build a lookup of section index → section name (from content slides)
    sec_idx_to_name: dict[int, str] = {}
    for slide_obj in slides:
        si = getattr(slide_obj, "section_index", -1)
        sn = getattr(slide_obj, "section", None) or ""
        if si >= 0 and sn and si not in sec_idx_to_name:
            sec_idx_to_name[si] = sn

    for pptx_idx, slide_obj in enumerate(slides):
        is_agenda = getattr(slide_obj, "is_agenda_slide", False)
        hl_idx    = getattr(slide_obj, "agenda_highlight_index", -1)

        if is_agenda and hl_idx >= 0:
            # Highlighted agenda → belongs to the section it introduces
            sec_name = sec_idx_to_name.get(hl_idx, f"Section {hl_idx + 1}")
        elif is_agenda:
            # Overview agenda (hl=-1) → put in the first section (cover)
            sec_name = sec_idx_to_name.get(0, "Cover")
        else:
            sec_name = getattr(slide_obj, "section", None) or sec_idx_to_name.get(0, "Cover")

        if sec_name not in section_map:
            section_map[sec_name] = []
            section_order.append(sec_name)
        section_map[sec_name].append(pptx_idx)

    if len(section_order) < 2:
        return   # nothing meaningful to inject

    # Build rId lookup: pptx slide index → relationship id
    def _rId(pptx_slide_index: int) -> str | None:
        target_part = prs.slides[pptx_slide_index].part
        for rId, rel in prs.part.rels.items():
            if hasattr(rel, "_target") and rel._target is target_part:
                return rId
        return None

    P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    sec_lst_el = etree.Element(f"{{{P14}}}sectionLst",
                               nsmap={"p14": P14, "r": R_NS})

    for sec_name in section_order:
        sec_el = etree.SubElement(sec_lst_el, f"{{{P14}}}section")
        sec_el.set("name", sec_name or "Default")
        sec_el.set("id", "{" + str(_uuid.uuid4()).upper() + "}")
        sld_id_lst = etree.SubElement(sec_el, f"{{{P14}}}sldIdLst")
        for pptx_idx in section_map.get(sec_name, []):
            rId = _rId(pptx_idx)
            if rId:
                sld_el = etree.SubElement(sld_id_lst, f"{{{P14}}}sldId")
                sld_el.set(f"{{{R_NS}}}id", rId)

    # Inject into presentation extLst
    prs_el = prs.element
    extLst = prs_el.find(qn("p:extLst"))
    if extLst is None:
        extLst = etree.SubElement(prs_el, qn("p:extLst"))

    # Remove any existing section extension to avoid duplicates
    for ext in list(extLst):
        if ext.get("uri") == "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}":
            extLst.remove(ext)

    ext_el = etree.SubElement(extLst, qn("p:ext"))
    ext_el.set("uri", "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}")
    ext_el.append(sec_lst_el)


# ── PptxBuilder ───────────────────────────────────────────────────────────────

class PptxBuilder:
    """Thin orchestrator: scaffolds PPTX slides, delegates layout to TEMPLATE_REGISTRY."""

    def __init__(self, ds: DesignSystem) -> None:
        self.ds = ds

    # -- Public API ------------------------------------------------------------

    def build(self, slides: list, output_path: Path) -> None:
        """Build a complete .pptx file from *slides* and write to *output_path*."""
        ds     = self.ds
        width, height, _ = self._canvas()
        prs    = Presentation()
        prs.slide_width  = Emu(px(width))
        prs.slide_height = Emu(px(height))
        total = len(slides)
        print(f"Building PPTX output: {total} slide(s) -> {output_path}")
        for i, slide_obj in enumerate(slides):
            self._build_slide(prs, slide_obj, page_num=i + 1, total=total)
            print(f"  Slide {i + 1}: {slide_obj.title}")
        output_path.parent.mkdir(parents=True, exist_ok=True)
        _inject_pptx_sections(prs, slides)
        prs.save(str(output_path))
        print(f"Wrote {output_path}")

    # -- Per-slide rendering ---------------------------------------------------

    def _build_slide(self, prs, slide: Slide, page_num: int = 1, total: int = 1) -> None:
        """Append a single slide to *prs*."""
        ds              = self.ds
        width, height, margin = self._canvas()
        chrome          = self._resolve_chrome(slide)

        pptx_slide = self._make_slide(prs)
        ctx        = PptxCtx(pptx_slide, ds)
        template   = slide.template_hint or "hero-title"

        self._background(ctx, chrome, width, height)

        # Optional accent bar at very top of slide
        accent_offset = 0
        if chrome["accent_show"]:
            self._accent_bar(ctx, chrome, width)
            accent_offset = chrome["accent_height"]

        # Logos in the header area
        self._logos(ctx, chrome, width, margin, accent_offset)

        # Slide title → returns content_y (first free Y below title + gap)
        title_y   = margin + accent_offset
        content_y = self._slide_title(ctx, slide.title, chrome,
                                      margin, title_y, width - 2 * margin)

        # Header divider line (width + alignment driven by CSS tokens)
        if chrome["divider_show"]:
            dx, dw = self._divider_span(width, margin, chrome["divider_width"],
                                        chrome["divider_align"])
            ctx.divider(dx, content_y - 4, dw, color=chrome["divider_color"])

        # Footer zone (rendered below content; reduces content_h)
        footer_h = 0
        if chrome["footer_show"] or chrome["page_number_show"]:
            footer_h = chrome["footer_height"]
            self._footer(ctx, chrome, slide, width, height, margin, page_num, total)

        content_h = height - content_y - margin - footer_h

        # Content-area background (optional fill between header and footer)
        self._content_bg(ctx, chrome, margin, content_y, width, content_h)

        blocks = _extract_section_blocks(slide)
        layout = TEMPLATE_REGISTRY.get(template, TEMPLATE_REGISTRY["hero-title"])
        layout.render(ctx, slide, blocks, margin, content_y, content_h,
                      width, height, dispatch)

    # -- Scaffold helpers ──────────────────────────────────────────────────────

    def _background(self, ctx: PptxCtx, chrome: dict, width: int, height: int) -> None:
        """Fill the whole slide; uses --slide-bg-color (falls back to 'surface')."""
        bg = chrome.get("slide_bg_color") or ctx.color("surface")
        ctx.rect(0, 0, width, height, fill=bg)

    def _slide_title(self, ctx: PptxCtx, title: str, chrome: dict,
                     x: int, y: int, w: int) -> int:
        """Render slide title; return content_y (Y after title + gap)."""
        h = chrome["header_height"]
        ctx.text(x, y, w, h, title,
                 size=ctx.font_size("heading"), bold=ctx.font_bold("heading"),
                 color=chrome["title_color"],
                 align=chrome["title_align"], valign="top")
        return y + h + chrome["header_gap"]

    def _make_slide(self, prs):
        """Create a blank slide (no placeholders) and return it."""
        blank_layout = prs.slide_layouts[6]  # Blank
        return prs.slides.add_slide(blank_layout)

    def _canvas(self):
        c = self.ds.canvas()
        return (int(c.get("width", 1920)),
                int(c.get("height", 1080)),
                int(c.get("margin", 80)))

    def _accent_bar(self, ctx: PptxCtx, chrome: dict, width: int) -> None:
        """Draw a thin accent stripe across the very top of the slide."""
        ctx.rect(0, 0, width, chrome["accent_height"],
                 fill=chrome["accent_color"], stroke=None)

    def _content_bg(self, ctx: PptxCtx, chrome: dict,
                    margin: int, content_y: int, width: int, content_h: int) -> None:
        """Draw an optional background fill for the content area."""
        bg = chrome.get("content_area_bg_color", "transparent")
        if not ctx.is_transparent_fill(bg):
            ctx.rect(margin, content_y, width - 2 * margin, content_h, fill=bg)

    @staticmethod
    def _divider_span(slide_w: int, margin: int, width_token: str,
                      align_token: str) -> tuple[int, int]:
        """Compute (x, w) for a header or footer divider from CSS width/align tokens.

        *width_token*  e.g. ``"100%"``, ``"50%"`` or ``"480"`` (px).
        *align_token*  one of ``"left"`` | ``"center"`` | ``"right"``.
        Returns ``(start_x, span_w)``.
        """
        inner = slide_w - 2 * margin
        raw   = str(width_token).strip().lower()
        try:
            if raw.endswith("%"):
                ratio  = max(0.0, min(1.0, float(raw[:-1]) / 100.0))
                span_w = max(1, int(inner * ratio))
            else:
                val = float(raw)
                span_w = max(1, int(val if val <= inner else inner))
        except (ValueError, TypeError):
            span_w = inner
        align = str(align_token).strip().lower()
        if align == "center":
            start_x = margin + max(0, (inner - span_w) // 2)
        elif align == "right":
            start_x = margin + max(0, inner - span_w)
        else:
            start_x = margin
        return start_x, span_w

    def _logos(self, ctx: PptxCtx, chrome: dict, width: int,
               margin: int, accent_offset: int = 0) -> None:
        """Render primary (right) and secondary (left) logos in the header area."""
        logo_y = margin + accent_offset
        for side in ("primary", "secondary"):
            src = chrome.get(f"logo_{side}_src", "none")
            if not src or str(src).strip().lower() in ("none", "", "false", "0"):
                continue
            lw = chrome.get(f"logo_{side}_width",  120 if side == "primary" else 80)
            lh = chrome.get(f"logo_{side}_height", 40  if side == "primary" else 30)
            # Resolve file path relative to the CSS source directory
            from pathlib import Path
            from io import BytesIO
            from pptx.util import Emu
            src_path = Path(src)
            if not src_path.is_absolute():
                css_dir = getattr(self.ds._css, "source_dir", None)
                if css_dir:
                    src_path = css_dir / src_path
            if not src_path.exists():
                continue
            lx = (width - margin - lw if side == "primary"
                  else margin)
            ly = logo_y + max(0, (chrome["header_height"] - lh) // 2)
            try:
                ctx.slide.shapes.add_picture(
                    BytesIO(src_path.read_bytes()),
                    Emu(ctx._px(lx)), Emu(ctx._px(ly)),
                    Emu(ctx._px(lw)), Emu(ctx._px(lh)),
                )
            except Exception:
                pass  # skip gracefully on invalid image

    def _footer(self, ctx: PptxCtx, chrome: dict, slide: Slide,
                width: int, height: int, margin: int,
                page_num: int, total: int) -> None:
        """Draw footer zone: optional divider, left text, right page number."""
        fh      = chrome["footer_height"]
        zone_y  = height - margin - fh
        text_y  = zone_y + 4
        text_h  = fh - 8
        font_sz = max(10, ctx.font_size("body") - 4)
        inner_w = width - 2 * margin

        if chrome["footer_divider_show"]:
            fdx, fdw = self._divider_span(width, margin,
                                          chrome.get("footer_divider_width", "100%"),
                                          chrome.get("footer_divider_align", "left"))
            ctx.divider(fdx, zone_y, fdw, color=chrome["footer_divider_color"])

        if chrome["page_number_show"]:
            num_text = self._format_page_number(page_num, total,
                                                chrome["page_number_format"])
            ctx.text(margin, text_y, inner_w, text_h, num_text,
                     size=font_sz, bold=False,
                     color=chrome["page_number_color"],
                     align=chrome["page_number_align"], valign="middle")

        footer_text = self._resolve_footer_text(chrome, slide)
        if footer_text and chrome["footer_show"]:
            text_w = (
                (inner_w * 2) // 3
                if chrome["page_number_show"] and chrome["page_number_align"] == "right"
                else inner_w
            )
            ctx.text(margin, text_y, text_w, text_h, footer_text,
                     size=font_sz, bold=False,
                     color=chrome["footer_text_color"],
                     align=chrome["footer_text_align"], valign="middle")

    @staticmethod
    def _format_page_number(page_num: int, total: int, fmt: str) -> str:
        fmt = str(fmt).strip()
        if "/" in fmt or "total" in fmt.lower():
            return f"{page_num} / {total}"
        return str(page_num)

    def _resolve_footer_text(self, chrome: dict, slide: Slide) -> str:
        """Footer text: per-slide override > deck front-matter > CSS token."""
        if slide.chrome_overrides.get("footer_text") is not None:
            return str(slide.chrome_overrides["footer_text"])
        if slide.front_matter.get("footer_text"):
            return str(slide.front_matter["footer_text"])
        return chrome.get("footer_text", "")

    def _resolve_chrome(self, slide: Slide) -> dict:
        """Merge CSS chrome defaults with deck front-matter and per-slide overrides."""
        chrome = dict(self.ds.chrome())
        fm = slide.front_matter or {}
        if fm.get("page_number_format"):
            chrome["page_number_format"] = str(fm["page_number_format"])
        if "show_page_numbers" in fm:
            chrome["page_number_show"] = bool(fm["show_page_numbers"])
        if "show_footer" in fm:
            chrome["footer_show"] = bool(fm["show_footer"])
        if fm.get("footer_text"):
            chrome["footer_text"] = str(fm["footer_text"])
        co = slide.chrome_overrides or {}
        for key in ("accent_show", "footer_show", "page_number_show",
                    "divider_show", "footer_divider_show"):
            if key in co:
                chrome[key] = bool(co[key])
        if co.get("footer_text") is not None:
            chrome["footer_text"] = str(co["footer_text"])
        if co.get("title_align"):
            chrome["title_align"] = str(co["title_align"])
        return chrome


# ── Convenience function (backwards-compatible) ───────────────────────────────

def _resolve_stylesheet(md_path: Path, stylesheet_path: Path | None) -> Path | None:
    """Resolve the CSS stylesheet to use for *md_path*.

    Resolution order:
    1. Explicit *stylesheet_path* argument (highest priority).
    2. ``theme:`` key in the deck's YAML front-matter.
    3. ``theme.css`` alongside the Markdown file (auto-detect).
    4. Fall through to ``None`` → DesignSystem uses the bundled default.
    """
    if stylesheet_path is not None:
        return stylesheet_path

    # Check front-matter for a theme: path
    try:
        text = md_path.read_text(encoding="utf-8")
        import re as _re
        fm_match = _re.match(r"^---\n(.+?)\n---\n", text, _re.DOTALL)
        if fm_match:
            import yaml as _yaml
            fm = _yaml.safe_load(fm_match.group(1)) or {}
            theme_val = fm.get("theme")
            if theme_val:
                candidate = (md_path.parent / theme_val).resolve()
                if candidate.exists():
                    return candidate
    except Exception:
        pass

    # Auto-detect theme.css in the same directory as the deck
    sibling = md_path.parent / "theme.css"
    if sibling.exists():
        return sibling

    return None


def build_pptx(md_path: Path, stylesheet_path, output_path: Path, qa: bool = True) -> None:
    """Build a .pptx file from Markdown + CSS stylesheet."""
    css    = _resolve_stylesheet(md_path, Path(stylesheet_path) if stylesheet_path else None)
    if css:
        print(f"Using stylesheet: {css}")
    else:
        print("WARNING: no project theme.css found — falling back to materialdesign3 default",
              file=__import__('sys').stderr)
    ds     = DesignSystem.load(css)
    slides = parse_markdown(md_path)
    if slides and slides[0].front_matter.get("auto-agenda", True):
        slides = inject_agenda_slides(slides)
        strip_agenda_from_deck(md_path)   # keep source clean; slides are injected at build time
    PptxBuilder(ds).build(slides, output_path)

    if qa:
        try:
            qa_output_dir = output_path.parent / "qa"
            print(f"Running mandatory visual QA (slides -> {qa_output_dir})")
            qa_render(output_path, output_dir=qa_output_dir, dpi=150)
        except Exception as exc:
            print(f"WARNING: QA render failed: {exc}", file=__import__('sys').stderr)


# ── CLI entrypoint ────────────────────────────────────────────────────────────

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Build PPTX from Markdown")
    parser.add_argument("md_file",         help="Input Markdown file")
    parser.add_argument("--stylesheet",    help="CSS theme stylesheet path (theme.css)",
                        default=None)
    parser.add_argument("--output",        help="Output .pptx path",
                        default=None)
    parser.add_argument("--no-qa", action="store_true",
                        help="Skip mandatory visual QA stage (not recommended)")
    args = parser.parse_args()

    md  = Path(args.md_file)
    css = Path(args.stylesheet) if args.stylesheet else None
    out = Path(args.output) if args.output else md.with_suffix(".pptx")
    build_pptx(md, css, out, qa=not args.no_qa)
